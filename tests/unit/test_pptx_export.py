"""PPTX text extraction — slides to DocumentBlocks, pictures ignored."""
from __future__ import annotations

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from core.document import _build_md_from_blocks, _pptx_to_blocks  # noqa: E402


@pytest.fixture()
def deck(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Title A"
    tf = slide.placeholders[1].text_frame
    tf.text = "top"
    nested = tf.add_paragraph()
    nested.text = "child"
    nested.level = 1

    second = prs.slides.add_slide(prs.slide_layouts[5])
    second.shapes.title.text = "Title B"
    table = second.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "h1"
    table.cell(0, 1).text = "h2"
    table.cell(1, 0).text = "v1"
    table.cell(1, 1).text = "v2"

    png = tmp_path / "p.png"
    Image = pytest.importorskip("PIL.Image", reason="Pillow needed for picture")
    Image.new("RGB", (10, 10)).save(png)
    second.shapes.add_picture(str(png), Inches(1), Inches(4), Inches(1))

    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


def test_title_appears_once(deck) -> None:
    """python-pptx returns a fresh proxy per `.title` access, so an identity
    check against iterated shapes fails and the title is emitted twice."""
    blocks = _pptx_to_blocks(deck)
    titles = [b.text for b in blocks if b.text in ("Title A", "Title B")]
    assert titles == ["Title A", "Title B"], titles
    assert [b.block_type for b in blocks if b.text == "Title A"] == ["heading"]


def test_structure_and_nesting(deck) -> None:
    blocks = _pptx_to_blocks(deck)
    kinds = [b.block_type for b in blocks]
    assert kinds.count("heading") == 2
    assert "table" in kinds
    nested = [b for b in blocks if b.text == "child"]
    assert nested and nested[0].level == 2, "outline level must drive nesting"


def test_pptx_pictures_carry_no_bytes_so_nothing_is_written(deck, tmp_path) -> None:
    """PPTX reading records a picture's presence but not its bytes, so the
    Markdown writer has nothing to place — and must not invent a broken link."""
    out = tmp_path / "d.md"
    summary = _build_md_from_blocks(_pptx_to_blocks(deck), str(out))
    text = out.read_text(encoding="utf-8")
    assert "![" not in text
    assert not list(tmp_path.glob("*_images"))
    assert summary.images == 0


def test_empty_deck_does_not_crash(tmp_path) -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "empty.pptx"
    prs.save(str(path))
    assert _pptx_to_blocks(str(path)) == []
