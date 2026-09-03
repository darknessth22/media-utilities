"""PDF/DOCX -> PPTX: the deck must be valid AND readable.

"Valid" is checked structurally (OOXML parts, well-formed slide XML) because a
malformed package is what makes PowerPoint show its repair prompt. "Readable"
means no slide is overloaded or empty — a naive one-page-per-slide conversion
produces text running off the bottom, which looks broken even when the file is
technically fine.
"""
from __future__ import annotations

import os
import xml.etree.ElementTree as ET
import zipfile

import pytest

pytest.importorskip("pptx")
pytest.importorskip("docx")

from docx import Document  # noqa: E402
from pptx import Presentation  # noqa: E402

from core.document import (  # noqa: E402
    _SLIDE_MAX_CHARS, DocumentBlock, _build_pptx_from_blocks,
    _group_blocks_into_slides, convert_document,
)


def _blocks(n_sections=3, per_section=8):
    out = [DocumentBlock("heading", text="Report", level=1)]
    for s in range(n_sections):
        out.append(DocumentBlock("heading", text=f"Section {s}", level=2))
        for i in range(per_section):
            out.append(DocumentBlock(
                "list_item", text=f"point {i} in section {s} with some length",
                level=1, list_style="bullet"))
    return out


def _assert_valid_package(path: str) -> Presentation:
    z = zipfile.ZipFile(path)
    assert z.testzip() is None, "zip is corrupt"
    names = z.namelist()
    for required in ("[Content_Types].xml", "ppt/presentation.xml", "_rels/.rels"):
        assert required in names, f"missing {required}"
    for slide in (n for n in names if n.startswith("ppt/slides/slide")):
        ET.fromstring(z.read(slide))       # raises if malformed
    return Presentation(path)


def test_package_is_structurally_valid(tmp_path) -> None:
    out = tmp_path / "d.pptx"
    _build_pptx_from_blocks(_blocks(), str(out))
    prs = _assert_valid_package(str(out))
    assert len(prs.slides) > 0


def test_long_content_is_split_not_overflowed(tmp_path) -> None:
    out = tmp_path / "d.pptx"
    _build_pptx_from_blocks(_blocks(n_sections=4, per_section=20), str(out))
    prs = _assert_valid_package(str(out))
    for slide in prs.slides:
        chars = sum(len(sh.text_frame.text)
                    for sh in slide.shapes if sh.has_text_frame)
        assert chars < _SLIDE_MAX_CHARS + 300, f"slide overloaded: {chars} chars"


def test_no_empty_slides(tmp_path) -> None:
    """A trailing flush used to leave a title-only '(cont.)' slide behind."""
    blocks = _blocks(n_sections=2, per_section=6)
    blocks.append(DocumentBlock("table", table_data=[["a", "b"], ["1", "2"]]))
    out = tmp_path / "d.pptx"
    _build_pptx_from_blocks(blocks, str(out))
    prs = _assert_valid_package(str(out))
    for i, slide in enumerate(prs.slides, 1):
        has = any(
            (sh.has_text_frame and sh.text_frame.text.strip())
            or (getattr(sh, "has_table", False) and sh.has_table)
            for sh in slide.shapes)
        assert has, f"slide {i} is empty"


def test_shapes_stay_on_the_slide(tmp_path) -> None:
    blocks = _blocks(n_sections=2)
    blocks.append(DocumentBlock(
        "table", table_data=[[f"c{c}" for c in range(4)] for _ in range(30)]))
    out = tmp_path / "d.pptx"
    _build_pptx_from_blocks(blocks, str(out))
    prs = _assert_valid_package(str(out))
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if None in (shape.left, shape.top, shape.width, shape.height):
                continue
            assert shape.left >= 0 and shape.top >= 0, f"slide {i}: negative origin"
            assert shape.left + shape.width <= prs.slide_width + 1000
            assert shape.top + shape.height <= prs.slide_height + 1000


def test_tall_table_is_split_with_repeated_header(tmp_path) -> None:
    rows = [["h0", "h1"]] + [[f"r{r}", f"v{r}"] for r in range(40)]
    groups = _group_blocks_into_slides([DocumentBlock("table", table_data=rows)])
    assert len(groups) > 1, "a 41-row table must span several slides"
    for group in groups:
        assert group["blocks"][0].table_data[0] == ["h0", "h1"], "header not repeated"


def test_headings_start_new_slides(tmp_path) -> None:
    groups = _group_blocks_into_slides([
        DocumentBlock("heading", text="One", level=1),
        DocumentBlock("paragraph", text="a"),
        DocumentBlock("heading", text="Two", level=1),
        DocumentBlock("paragraph", text="b"),
    ])
    assert [g["title"] for g in groups] == ["One", "Two"]


def test_images_are_skipped_and_reported(tmp_path) -> None:
    blocks = [DocumentBlock("heading", text="H", level=1),
              DocumentBlock("paragraph", text="text"),
              DocumentBlock("image", image_bytes=b"x", image_ext="png", page_num=0)]
    out = tmp_path / "d.pptx"
    summary = _build_pptx_from_blocks(blocks, str(out))
    prs = _assert_valid_package(str(out))
    assert summary.skipped_elements
    for slide in prs.slides:
        assert not any(sh.shape_type == 13 for sh in slide.shapes)  # 13 = PICTURE


def test_empty_input_still_produces_a_readable_file(tmp_path) -> None:
    """An empty package can trip readers; leave one explanatory slide."""
    out = tmp_path / "d.pptx"
    _build_pptx_from_blocks([], str(out))
    prs = _assert_valid_package(str(out))
    assert len(prs.slides) == 1


def test_docx_to_pptx_end_to_end(tmp_path) -> None:
    doc = Document()
    doc.add_heading("Title", 1)
    doc.add_paragraph("Body text here.")
    for i in range(12):
        doc.add_paragraph(f"item {i}", style="List Bullet")
    table = doc.add_table(rows=3, cols=2)
    for r in range(3):
        for c in range(2):
            table.cell(r, c).text = f"{r}{c}"
    src = tmp_path / "in.docx"
    doc.save(str(src))

    ok, out, summary = convert_document(str(src), "pptx")
    assert ok, out
    prs = _assert_valid_package(out)
    assert len(prs.slides) >= 2
    assert summary.tables == 1


def test_pptx_output_rejected_for_unsupported_input(tmp_path) -> None:
    img = tmp_path / "x.png"
    pytest.importorskip("PIL.Image")
    from PIL import Image
    Image.new("RGB", (8, 8)).save(img)
    ok, msg, _ = convert_document(str(img), "pptx")
    assert not ok and "not supported" in msg.lower()


# ── design decks ──────────────────────────────────────────────────────────

def _deck_pdf(tmp_path, pages=4, size=(1920, 1080)):
    """A composed 16:9 deck: full-bleed background with text over it."""
    fitz = pytest.importorskip("fitz")
    Image = pytest.importorskip("PIL.Image")
    bg = tmp_path / "bg.png"
    Image.new("RGB", (1600, 900), (180, 205, 225)).save(bg)

    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page(width=size[0], height=size[1])
        # Oversized on purpose — real decks bleed past the page edge.
        page.insert_image(fitz.Rect(-100, -80, size[0] + 100, size[1] + 80),
                          filename=str(bg))
        page.insert_text((300, 300), f"SLIDE {i} HEADLINE", fontsize=48)
    path = tmp_path / "deck.pdf"
    doc.save(str(path))
    doc.close()
    return str(path)


def test_design_deck_is_detected(tmp_path) -> None:
    fitz = pytest.importorskip("fitz")
    from core.document import _pdf_is_design_deck, _page_is_designed_slide
    with fitz.open(_deck_pdf(tmp_path)) as doc:
        assert _pdf_is_design_deck(doc)
        assert _page_is_designed_slide(doc[0])


def test_portrait_document_is_not_a_deck(tmp_path) -> None:
    """A4 text pages must keep the text-rebuilding path."""
    fitz = pytest.importorskip("fitz")
    from core.document import _pdf_is_design_deck
    doc = fitz.open()
    for _ in range(3):
        page = doc.new_page()          # default A4-ish portrait
        page.insert_text((72, 100), "Heading", fontsize=20)
        page.insert_text((72, 140), "Body text " * 12, fontsize=11)
    path = tmp_path / "doc.pdf"
    doc.save(str(path))
    doc.close()
    with fitz.open(str(path)) as d2:
        assert not _pdf_is_design_deck(d2)


def test_deck_converts_one_slide_per_page(tmp_path) -> None:
    """The layout is the content, so a page must not be re-flowed into several
    slides — the sample deck produced 57 slides from 36 pages before this."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=5)
    ok, out, summary = convert_document(path, "pptx")
    assert ok, out
    prs = _assert_valid_package(out)
    assert len(prs.slides) == 5, f"expected 5 slides, got {len(prs.slides)}"
    assert summary.images == 5


def test_deck_slide_matches_source_aspect_and_fills_it(tmp_path) -> None:
    """A 16:9 page on a 4:3 slide gets cropped; adopt the source shape."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=2)
    ok, out, _ = convert_document(path, "pptx")
    assert ok
    prs = _assert_valid_package(out)
    ratio = prs.slide_width / prs.slide_height
    assert abs(ratio - 16 / 9) < 0.02, f"slide ratio {ratio:.3f}"
    pic = prs.slides[0].shapes[0]
    assert pic.left == 0 and pic.top == 0, "picture is offset, so it is cropped"
    assert abs(pic.width - prs.slide_width) < 2000
    assert abs(pic.height - prs.slide_height) < 2000


def test_deck_slides_carry_no_placeholder_prompts(tmp_path) -> None:
    """An empty 'Click to add title' box over the artwork is what the user
    reported; a full-bleed slide must use the blank layout."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=2)
    ok, out, _ = convert_document(path, "pptx")
    assert ok
    prs = _assert_valid_package(out)
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            assert not shape.is_placeholder, f"slide {i} has a placeholder"


def test_deck_output_stays_a_sane_size(tmp_path) -> None:
    """Rendering at full DPI as PNG made a 36-page deck 394 MB."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=6)
    ok, out, _ = convert_document(path, "pptx")
    assert ok
    mb = os.path.getsize(out) / 1_000_000
    assert mb < 6, f"{mb:.1f} MB for 6 slides is too large"


def test_deck_text_is_editable_not_baked_into_the_image(tmp_path) -> None:
    """The whole point: keep the design AND keep the wording editable."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=3)
    ok, out, summary = convert_document(path, "pptx")
    assert ok, out
    prs = _assert_valid_package(out)
    for i, slide in enumerate(prs.slides, 1):
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]
        boxes = [sh for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        assert pics, f"slide {i} lost its background"
        assert boxes, f"slide {i} has no editable text"
        assert "HEADLINE" in " ".join(b.text_frame.text for b in boxes)
    assert summary.text_blocks >= 3


def test_deck_background_has_the_text_removed(tmp_path) -> None:
    """Rendering the text into the background AND overlaying it would show the
    old glyphs underneath any edit."""
    fitz = pytest.importorskip("fitz")
    from core.document import _render_page_block
    path = _deck_pdf(tmp_path, pages=1)
    with fitz.open(path) as doc:
        page = doc[0]
        assert page.get_text().strip(), "fixture has no text to begin with"
        block = _render_page_block(page, 0, editable_text=True)
        # The page object is redacted in place before rendering.
        assert not page.get_text().strip(), "text still in the background"
    assert block is not None and block.text_overlays


def test_overlay_font_size_is_scaled_to_the_slide(tmp_path) -> None:
    """slide_height is EMU and the source size is points; dividing one by the
    other gave a 257,000 pt font."""
    from core.document import convert_document
    path = _deck_pdf(tmp_path, pages=1)
    ok, out, _ = convert_document(path, "pptx")
    assert ok
    prs = _assert_valid_package(out)
    sizes = [run.font.size.pt
             for slide in prs.slides for shape in slide.shapes
             if shape.has_text_frame
             for para in shape.text_frame.paragraphs for run in para.runs
             if run.font.size is not None]
    assert sizes, "no sized runs found"
    for pt in sizes:
        assert 4 <= pt <= 200, f"implausible font size {pt}pt"


def test_alignment_is_inferred_from_geometry(tmp_path) -> None:
    """A PDF stores no alignment; centred text has to be inferred from line
    geometry or every edit drifts sideways."""
    fitz = pytest.importorskip("fitz")
    from core.document import _page_text_overlays

    def build(name, centred):
        doc = fitz.open()
        page = doc.new_page(width=1920, height=1080)
        font = fitz.Font("helv")
        for text, y in (("SHORT LINE", 300), ("A LONGER LINE HERE", 348)):
            width = font.text_length(text, fontsize=40)
            x = (960 - width / 2) if centred else 700
            page.insert_text((x, y), text, fontsize=40)
        path = tmp_path / name
        doc.save(str(path))
        doc.close()
        with fitz.open(str(path)) as d2:
            return _page_text_overlays(d2[0])

    centred = build("centred.pdf", True)
    assert centred and any(sp["align"] == "center" for sp in centred), \
        f"centred not detected: {[sp['align'] for sp in centred]}"

    # Lines sharing a left edge must NOT be reported as centred.
    left = build("left.pdf", False)
    assert left and all(sp["align"] == "left" for sp in left), \
        f"left misread: {[sp['align'] for sp in left]}"


# ---------------------------------------------------------------------------
# Letter spacing, paragraph breaks and RTL (see bug-094)
# ---------------------------------------------------------------------------

def _chars(text, glyph_w, space_w, size):
    """Per-character boxes the way get_text('rawdict') reports them."""
    out, x = [], 0.0
    for ch in text:
        w = space_w if ch == " " else glyph_w
        out.append({"c": ch, "bbox": (x, 0.0, x + w, size)})
        x += w
    return out


def test_fake_tracking_restores_word_breaks():
    """"T O L E O  T O W E R" must come back as "TOLEO TOWER" with tracking.

    Design decks spell letter spacing with literal spaces — one between
    letters, two between words. Passed through verbatim, PowerPoint sets those
    at its own width and the line reads "TOLEOTOWER".
    """
    from core.document import _measure_tracking

    text, em = _measure_tracking(
        _chars("T O L E O  T O W E R", 36.0, 10.8, 65.0), 65.0)
    assert text == "TOLEO TOWER"
    assert em > 0.1


def test_real_tracking_is_measured_not_invented():
    from core.document import _measure_tracking

    wide = _measure_tracking(_chars("TOLEO", 20.0, 0.0, 40.0), 40.0)
    # No spaces at all, but the advances are wide: tracking must be reported.
    spaced = [{"c": c, "bbox": (i * 28.0, 0.0, i * 28.0 + 20.0, 40.0)}
              for i, c in enumerate("TOLEO")]
    text, em = _measure_tracking(spaced, 40.0)
    assert text == "TOLEO"
    assert 0.15 < em < 0.25
    # Ordinary prose claims none.
    text, em = _measure_tracking(_chars("Hello world today", 7.0, 7.0, 14.0), 14.0)
    assert text == "Hello world today"
    assert em == 0.0
    assert wide[0] == "TOLEO"


def test_rtl_text_is_returned_in_logical_order():
    """PDF stores Arabic visually reversed, in presentation forms.

    The fixture is derived from the expected string rather than typed out, so
    a typo in a presentation-form literal cannot make the test lie (it did on
    the first attempt — a stray waw for a noon).
    """
    import unicodedata
    from core.document import _fix_rtl_order

    logical = "مرحبا بكم "               "في زنجبار"
    visual = " ".join(w[::-1] for w in reversed(logical.split(" ")))
    assert _fix_rtl_order(visual) == logical
    # Presentation forms must normalise back too.
    forms = unicodedata.normalize("NFKC", visual)
    assert _fix_rtl_order(forms) == logical
    # Latin is left exactly as it is.
    assert _fix_rtl_order("TOLEO TOWER") == "TOLEO TOWER"


def test_paragraph_breaks_survive_the_line_merge(tmp_path):
    """A gap wider than line spacing must stay a paragraph break.

    Merging consecutive lines into one editable box is right, but the sample
    deck's body copy is three paragraphs: 0.89 em between them, -0.01 em
    within one. Without the distinction all three ran together.
    """
    import fitz
    from core.document import _page_text_overlays

    src = tmp_path / "paras.pdf"
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    # insert_text() places the BASELINE, so the visible gap is much smaller
    # than the y step: 48 pt here yields 0.912 em between paragraphs and
    # -0.231 em within one, matching the sample deck's 0.89 / -0.01. Anything
    # under ~42 pt leaves the paragraph gap negative and indistinguishable
    # from line spacing; much over it stops the lines merging at all (which is
    # also correct output, just not what this test is checking).
    y = 100.0
    for i in range(3):
        page.insert_text((80, y), "Paragraph %d line one" % i, fontsize=14)
        page.insert_text((80, y + 16), "and its second line", fontsize=14)
        y += 48
    doc.save(src)
    doc.close()

    doc = fitz.open(src)
    overlays = _page_text_overlays(doc[0])
    doc.close()

    lines = [l for o in overlays for l in o["lines"]]
    assert len(lines) == 6
    # Two of the six start a new paragraph (the first line of a box does not
    # carry the flag — there is nothing above it to space away from).
    assert sum(1 for l in lines if l.get("space_before")) == 2


# ---------------------------------------------------------------------------
# Drop shadows and decorative rules (bug-099, bug-100)
# ---------------------------------------------------------------------------

def test_headline_drop_shadow_image_is_removed(tmp_path):
    """A small image lying on the words is a text shadow, not artwork.

    The sample deck renders six headline shadows as separate soft-masked
    images. Redacting the text leaves the shadow behind as a grey smudge — the
    "shading behind the text" the design does not have.
    """
    import fitz
    from core.document import _render_page_background

    src = tmp_path / "shadow.pdf"
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    # Full-page background: big, and barely overlaps the text.
    bg = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 960, 540))
    bg.set_rect(bg.irect, (120, 140, 160))
    page.insert_image(page.rect, pixmap=bg)
    # The headline, and a small "shadow" image right on top of it.
    page.insert_text((100, 200), "A REFINED EXPRESSION", fontsize=40)
    shadow = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 64, 16))
    shadow.set_rect(shadow.irect, (30, 30, 30))
    page.insert_image(fitz.Rect(95, 165, 560, 210), pixmap=shadow)
    doc.save(src)
    doc.close()

    doc = fitz.open(src)
    page = doc[0]
    area = fitz.Rect(95, 165, 560, 210)

    def darkest():
        """Darkest pixel in the headline area — the shadow is near-black."""
        pix = page.get_pixmap(clip=area, alpha=False)
        return min(pix.pixel(x, y)[0]
                   for x in range(0, pix.width, 7)
                   for y in range(0, pix.height, 3))

    assert len(page.get_images(full=True)) == 2, "fixture: background + shadow"
    assert darkest() < 40, "fixture: the shadow should be dark"

    _render_page_background(page, 0)
    after = darkest()
    remaining = len(page.get_images(full=True))
    doc.close()

    # Assert on PIXELS, not on the image list: delete_image blanks the XObject
    # rather than unlinking it, so get_images() still reports two entries even
    # though nothing of the shadow is drawn any more.
    assert after >= 100, "shadow (and text) should be gone from the area"
    assert remaining == 2


def test_rule_inside_a_text_block_survives(tmp_path):
    """Vector art a text bbox happens to enclose must survive the redaction.

    Page 1 of the sample deck merges "PROPRITIES" and "TOLEO TOWER" into one
    block spanning y 614-714, so the rule at y 707-711 sits inside it. This is
    the general guarantee that removing text never removes design elements.
    """
    import fitz
    from core.document import _render_page_background

    src = tmp_path / "rule.pdf"
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    # The second line has to be wide enough that its block spans the rule —
    # in the real deck the block is 601-1285 pt because it covers both lines.
    page.insert_text((100, 200), "TOLEO TOWER", fontsize=40)
    page.insert_text((100, 224), "PROPRITIES AND MORE TEXT TO WIDEN THE BOX",
                     fontsize=10)
    rule = fitz.Rect(100, 214, 420, 217)
    page.draw_rect(rule, color=(0.3, 0.1, 0.1), fill=(0.3, 0.1, 0.1), width=0)
    doc.save(src)
    doc.close()

    doc = fitz.open(src)
    page = doc[0]
    blocks = [fitz.Rect(b["bbox"]) for b in page.get_text("dict")["blocks"]
              if b.get("type") == 0]
    # The rule really is enclosed by a text block, or the test proves nothing.
    assert any((rule & b).get_area() > rule.get_area() * 0.5 for b in blocks)
    _render_page_background(page, 0)
    # Text gone, rule still drawn.
    assert "TOLEO" not in page.get_text()
    survived = [d for d in page.get_drawings()
                if not (fitz.Rect(d["rect"]) & rule).is_empty]
    doc.close()
    assert survived, "the rule must be re-drawn after redaction"


# ---------------------------------------------------------------------------
# Design vector art must survive text removal (bug-102)
# ---------------------------------------------------------------------------

def test_colour_fields_and_swatches_survive(tmp_path):
    """Removing text must not remove the page's colour.

    apply_redactions(graphics=2) drops any art OVERLAPPING a redaction
    rectangle, so a full-page colour field that clips a text block by 4%, or a
    palette swatch beside a caption, vanished with the lettering — a colour
    system slide converted to blank white.
    """
    import fitz
    from core.document import _render_page_background

    src = tmp_path / "swatches.pdf"
    doc = fitz.open()
    page = doc.new_page(width=1920, height=1080)
    page.draw_rect(page.rect, color=(0.77, 0.76, 0.75),
                   fill=(0.77, 0.76, 0.75), width=0)
    swatch = fitz.Rect(65, 352, 550, 609)
    page.draw_rect(swatch, color=(0.07, 0.11, 0.18),
                   fill=(0.07, 0.11, 0.18), width=0)
    # A caption inside the swatch, so its block overlaps the artwork.
    page.insert_text((90, 420), "Primary Color", fontsize=34)
    doc.save(src)
    doc.close()

    doc = fitz.open(src)
    page = doc[0]

    def colour_at(x, y):
        pix = page.get_pixmap(clip=fitz.Rect(x, y, x + 2, y + 2), alpha=False)
        return pix.pixel(0, 0)[:3]

    assert colour_at(1500, 100)[0] > 180, "fixture: page field is light"
    assert colour_at(300, 560)[0] < 60, "fixture: swatch is dark"

    _render_page_background(page, 0)
    field_after = colour_at(1500, 100)
    swatch_after = colour_at(300, 560)
    text_after = page.get_text()
    doc.close()

    assert "Primary" not in text_after, "the text should still be removed"
    assert field_after[0] > 180, "page colour field must survive"
    assert swatch_after[0] < 60, "swatch must survive"


def test_overlay_keeps_the_pdf_top_left_corner():
    """Padding must not displace a box — it shifts small type off its line.

    Padding used to be a flat fraction of the SLIDE (1.2% of width, 2% of
    height), so every box moved by the same 23 pt left and 11.9 pt up whatever
    the type size: negligible under a 58 pt headline, but more than a whole
    line-height under 10 pt body copy, which is why small text sat above its
    row. Slack now only grows the box right and down.
    """
    from pptx import Presentation
    from pptx.util import Emu
    from core.document import _add_text_overlay

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = {
        "align": "left", "x": 0.25, "y": 0.40, "w": 0.2, "h": 0.03,
        "page_h": 595.0,
        "lines": [{"runs": [{"text": "Business Setup", "size": 10.1,
                             "track": 0.0, "color": 0x000000,
                             "bold": False, "italic": False}],
                   "space_before": False}],
    }
    _add_text_overlay(prs, slide, spec)

    box = [s for s in slide.shapes if s.has_text_frame][0]
    want_left = int(0.25 * prs.slide_width)
    want_top = int(0.40 * prs.slide_height)
    # Within a point of the reported corner, not tens of points above/left.
    assert abs(box.left - want_left) < 12700
    assert abs(box.top - want_top) < 12700
    # And the box is still at least as large as the text it holds.
    assert box.width >= int(0.2 * prs.slide_width)
    assert box.height >= int(0.03 * prs.slide_height)
