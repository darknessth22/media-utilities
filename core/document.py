"""Document conversion engine: PDF, DOCX, and image → PDF."""
import contextlib
import io
import os
import re
import tempfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from io import BytesIO

import fitz  # PyMuPDF
from docx import Document
from docx.shared import Inches, Pt
from PIL import Image



@dataclass
class SpanInfo:
    """Formatting metadata for a text span within a DocumentBlock."""
    text: str
    bold: bool
    italic: bool
    font_size: float
    font_name: str
    color: tuple[int, int, int]  # RGB color tuple (0-255 each)


@dataclass
class DocumentBlock:
    """A structural unit extracted from a PDF page during conversion."""
    block_type: str  # heading, paragraph, list_item, table, image, scanned_page
    text: str | None = None
    spans: list[SpanInfo] = field(default_factory=list)
    level: int | None = None  # Heading level (1-4) or list nesting depth
    list_style: str | None = None  # bullet or number
    table_data: list[list[str]] | None = None
    image_bytes: bytes | None = None
    image_ext: str | None = None
    bbox: tuple[float, float, float, float] = (0, 0, 0, 0)
    page_num: int = 0
    # Positioned text runs laid over a rendered slide background, so a design
    # deck keeps its layout while its wording stays editable.
    text_overlays: list = field(default_factory=list)


@dataclass
class ConversionSummary:
    """A record of elements processed during a document conversion operation."""
    total_pages: int = 0
    text_blocks: int = 0
    headings: int = 0
    tables: int = 0
    images: int = 0
    list_items: int = 0
    scanned_pages: int = 0
    skipped_elements: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


_IMAGE_EXTS = frozenset({".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"})
_VALID_DOCS  = frozenset({".pdf", ".docx", ".pptx"})
# Markdown is an OUTPUT format only, by design. The GUI has a Markdown
# editor/viewer for reading .md files (Qt renders it natively), but
# converting FROM markdown into PDF/DOCX/PPTX is deliberately not offered.
_VALID_OUTPUTS = frozenset({"pdf", "docx", "md", "pptx"})


def _image_bbox(page, img):
    """Bounding box of *img* on *page*, or None when it cannot be resolved.

    `get_image_bbox()` REQUIRES an entry from `get_images(full=True)`: the
    9-tuple a bare `get_images()` returns raises "need item of full page image
    list", which failed EVERY conversion of a PDF containing an image. It also
    raises for an image that is referenced but never actually placed on the
    page, so callers must tolerate None rather than assume a rectangle.
    """
    try:
        return page.get_image_bbox(img)
    except Exception:
        return None


# Vector-figure extraction. A chart is not one shape: a single bar chart came
# out of get_drawings() as 13 separate primitives (14 lines + 6 rects), so
# mapping each to a PowerPoint autoshape yields scattered lines and boxes rather
# than a figure. Clustered primitives are therefore grouped by proximity and the
# region is rasterised once, which is pixel-faithful and ~7 KB for a chart.
_FIGURE_DPI = 150
# Measured on a bar chart: bars sat 25 pt apart, so 18 pt left every bar its own
# cluster. 34 pt groups a chart without swallowing a neighbouring paragraph
# (body text blocks were >45 pt away in the same samples).
_FIGURE_GAP = 34          # pt; primitives closer than this belong to one figure
_FIGURE_MIN_SIDE = 40     # pt; smaller clusters are rules/underlines, not figures
_FIGURE_MIN_ITEMS = 2     # a lone rule or box is not a figure
_FIGURE_MAX_PAGE_FRAC = 0.92   # a cluster this big is the page itself


def _cluster_rects(rects, gap: float):
    """Merge rectangles into clusters, joining any pair within *gap* points.

    Union-find would be tidier but this runs on a handful of primitives per
    page; repeatedly absorbing overlaps is simpler to read and fast enough.
    """
    import fitz

    clusters: list = []
    for r in rects:
        box = fitz.Rect(r)
        # NOT `is_empty`: a horizontal rule has height 0 and a vertical one
        # width 0, and those axis/gridline primitives are precisely what joins
        # a chart's bars into one figure. Only reject unusable geometry.
        if box.is_infinite or (box.width <= 0 and box.height <= 0):
            continue
        box.normalize()
        grown = fitz.Rect(box.x0 - gap, box.y0 - gap, box.x1 + gap, box.y1 + gap)
        merged = [box]
        rest = []
        for existing in clusters:
            if grown.intersects(existing):
                merged.append(existing)
            else:
                rest.append(existing)
        combined = merged[0]
        for m in merged[1:]:
            combined |= m
        rest.append(combined)
        clusters = rest

    # One pass can leave two clusters that only became adjacent late; settle.
    changed = True
    while changed:
        changed = False
        out: list = []
        for c in clusters:
            grown = fitz.Rect(c.x0 - gap, c.y0 - gap, c.x1 + gap, c.y1 + gap)
            hit = None
            for i, o in enumerate(out):
                if grown.intersects(o):
                    hit = i
                    break
            if hit is None:
                out.append(c)
            else:
                out[hit] |= c
                changed = True
        clusters = out
    return clusters


def _extract_vector_figures(page, page_num: int) -> list[DocumentBlock]:
    """Group vector drawings into figures and rasterise each as a PNG block.

    Returns image blocks, so every existing serialiser handles them with no
    special case. Text is deliberately not excluded from the clip: a chart's
    own axis labels belong in the picture.
    """
    import fitz

    try:
        drawings = page.get_drawings()
    except Exception as exc:
        _log("get_drawings failed on page %d: %s" % (page_num + 1, exc))
        return []
    if not drawings:
        return []

    rects = []
    for d in drawings:
        r = d.get("rect")
        if r is not None:
            rects.append(fitz.Rect(r))
    if not rects:
        return []

    page_area = page.rect.get_area() or 1
    blocks: list[DocumentBlock] = []

    for cluster in _cluster_rects(rects, _FIGURE_GAP):
        # `intersects` is False for a zero-area line lying inside the cluster,
        # so count by containment of the primitive's corners instead.
        members = sum(1 for r in rects
                      if cluster.x0 - 1 <= r.x0 and r.x1 <= cluster.x1 + 1
                      and cluster.y0 - 1 <= r.y0 and r.y1 <= cluster.y1 + 1)
        if members < _FIGURE_MIN_ITEMS:
            continue
        # Both sides must be substantial: a long thin band is a rule, not a
        # figure, however wide it is.
        if cluster.width < _FIGURE_MIN_SIDE or cluster.height < _FIGURE_MIN_SIDE:
            continue
        if cluster.get_area() / page_area > _FIGURE_MAX_PAGE_FRAC:
            # The whole page is vector art (a full-page diagram or a border);
            # the scanned-page path already covers that case better.
            continue

        clip = fitz.Rect(cluster) + (-6, -6, 6, 6)
        clip &= page.rect          # never clip outside the page
        if clip.is_empty or clip.width < 4 or clip.height < 4:
            continue
        try:
            pix = page.get_pixmap(clip=clip, dpi=_FIGURE_DPI)
            data = pix.tobytes("png")
        except Exception as exc:
            _log("figure raster failed on page %d: %s" % (page_num + 1, exc))
            continue
        if not data:
            continue
        blocks.append(DocumentBlock(
            block_type="image",
            image_bytes=data,
            image_ext="png",
            bbox=tuple(clip),
            page_num=page_num,
        ))
    return blocks


def _is_scanned_page(page) -> bool:
    """Detect image-only pages: little text and large images."""
    text_len = len(page.get_text().strip())
    if text_len > 200:
        return False
        
    page_area = page.rect.width * page.rect.height
    if page_area <= 0:
        return False
    image_area = 0
    for img in page.get_images(full=True):
        bbox = _image_bbox(page, img)
        if bbox is None:
            continue
        image_area += (bbox.x1 - bbox.x0) * (bbox.y1 - bbox.y0)
        
    return image_area > 0.8 * page_area


def _extract_page_blocks(page, page_num: int, body_font_size: float, doc, extracted_images: set, warnings: list = None, whole_slide: bool = False) -> list[DocumentBlock]:
    """Extract structured blocks (text, tables, images) from a PDF page.

    whole_slide: render the page as one image instead of decomposing it. Used
    for design-deck pages, where text sits on a full-bleed background and
    re-flowing the pieces destroys the layout.
    """
    blocks = []
    import re

    if whole_slide:
        rendered = _render_page_block(page, page_num)
        if rendered is not None:
            return [rendered]
    
    # Handle scanned page detection (T016)
    if _is_scanned_page(page):
        images = page.get_images(full=True)
        placed = [(img, _image_bbox(page, img)) for img in images]
        placed = [(img, bb) for img, bb in placed if bb is not None]
        if placed:
            largest_img = max(placed, key=lambda pair: pair[1].get_area())[0]
            xref = largest_img[0]
            img_data = _safe_pixmap_png(doc, xref)
            if img_data:
                blocks.append(DocumentBlock(
                    block_type="scanned_page",
                    image_bytes=img_data,
                    image_ext="png",
                    bbox=tuple(page.rect),
                    page_num=page_num
                ))
                extracted_images.add(xref)
                return blocks

    table_bboxes = []
    try:
        tables = page.find_tables()
        for tab in tables.tables:
            try:
                bbox = tuple(tab.bbox)
                rows = tab.extract()
                if rows:
                    blocks.append(DocumentBlock(
                        block_type="table",
                        table_data=rows,
                        bbox=bbox,
                        page_num=page_num
                    ))
                    table_bboxes.append(bbox)
            except Exception as e:
                _log(f"Table extraction failed on page {page_num+1}: {e}")
                if warnings is not None:
                    warnings.append(f"Table extraction failed on page {page_num+1} (falling back to text)")
    except Exception as e:
        _log(f"find_tables failed on page {page_num+1}: {e}")
        if warnings is not None:
            warnings.append(f"Table detection failed on page {page_num+1} (falling back to text)")

    # Vector figures: charts and diagrams drawn as many primitives, grouped
    # and rasterised so they arrive as one picture instead of scattered lines.
    blocks.extend(_extract_vector_figures(page, page_num))

    # Image extraction pass (T016 requirement: catch all images, not just text-flow)
    for img in page.get_images(full=True):
        xref = img[0]
        if xref in extracted_images:
            continue
            
        bbox = _image_bbox(page, img)
        # An image can be referenced without ever being placed.
        if bbox is None:
            continue
        # Skip very small images (likely icons or decorative elements)
        if bbox.width < 10 or bbox.height < 10:
            continue
            
        img_data = _safe_pixmap_png(doc, xref)
        if img_data:
            blocks.append(DocumentBlock(
                block_type="image",
                image_bytes=img_data,
                image_ext="png",
                bbox=tuple(bbox),
                page_num=page_num
            ))
            extracted_images.add(xref)

    # Text extraction with formatting (T014)
    text_dict = page.get_text("dict", flags=7)
    for b in text_dict.get("blocks", []):
        if "lines" not in b:
            continue
            
        block_rect = fitz.Rect(b["bbox"])
        if any(block_rect.intersects(fitz.Rect(tbb)) for tbb in table_bboxes):
            continue
            
        block_text = ""
        spans = []
        max_size = 0
        is_bold = False
        
        # Collect text and spans first
        for line in b["lines"]:
            for s in line.get("spans", []):
                span_text = s["text"]
                if not span_text: # Keep empty spans but check text
                    continue
                
                # Decompose sRGB color
                color_int = s["color"]
                rgb = ((color_int >> 16) & 0xFF, (color_int >> 8) & 0xFF, color_int & 0xFF)
                
                bold = bool(s["flags"] & 16)
                italic = bool(s["flags"] & 1)
                
                font_name = s.get("font", "").lower()
                if not bold and ("bold" in font_name or "black" in font_name):
                    bold = True
                if not italic and ("italic" in font_name or "oblique" in font_name):
                    italic = True
                
                if s["size"] > max_size:
                    max_size = s["size"]
                    is_bold = bold
                
                spans.append(SpanInfo(
                    text=span_text,
                    bold=bold,
                    italic=italic,
                    font_size=round(s["size"], 1),
                    font_name=s.get("font", "Helvetica"),
                    color=rgb
                ))
                block_text += span_text

        if not block_text.strip():
            continue
            
        block_type = "paragraph"
        level = None
        list_style = None
        
        if max_size > 1.5 * body_font_size:
            block_type = "heading"
            level = 1
        elif max_size > 1.2 * body_font_size:
            block_type = "heading"
            level = 2
        elif is_bold and max_size > 1.05 * body_font_size:
            block_type = "heading"
            level = 3
            
        if block_type == "paragraph" and spans:
            is_list, style = _detect_list_item(block_text, b["lines"][0]["spans"][0])
            if is_list:
                block_type = "list_item"
                list_style = style
                level = 1
                
                # Strip leading marker to avoid double bullets (Critical #1)
                marker_regex = r"^(\d+|[a-zA-Z]|[ivxIVX]+)[.)]\s+|^\W\s+"
                cleaned_text = re.sub(marker_regex, "", block_text, count=1)
                if cleaned_text != block_text:
                    # Update first span text to remove the marker
                    marker_len = len(block_text) - len(cleaned_text)
                    if spans and len(spans[0].text) >= marker_len:
                        spans[0].text = spans[0].text[marker_len:].lstrip()
                    block_text = cleaned_text

        blocks.append(DocumentBlock(
            block_type=block_type,
            text=block_text.strip(),
            spans=spans,
            level=level,
            list_style=list_style,
            bbox=tuple(b["bbox"]),
            page_num=page_num
        ))
        
    return sorted(blocks, key=lambda x: (x.bbox[1], x.bbox[0]))


def _log(msg: str) -> None:
    print(f"[DocConvert] {msg}", flush=True)


def _detect_body_font_size(doc) -> float:
    """Statistical analysis of font sizes across document to find the body text baseline."""
    size_counts = {}
    for page in doc:
        blocks = page.get_text("dict").get("blocks", [])
        for b in blocks:
            if "lines" not in b:
                continue
            for l in b["lines"]:
                for s in l["spans"]:
                    size = round(s["size"], 1)
                    size_counts[size] = size_counts.get(size, 0) + len(s["text"])
    
    if not size_counts:
        return 12.0
        
    # Return the most frequent font size by character count
    return max(size_counts, key=size_counts.get)


def _detect_list_item(block_text: str, first_span: dict) -> tuple[bool, str | None]:
    """Heuristic list detection using block text and font metadata."""
    import re
    text = block_text.strip()
    if not text:
        return False, None
        
    # Check for bullet symbols in font name or specific unicode characters
    font_name = first_span.get("font", "").lower()
    is_bullet_font = any(x in font_name for x in ["symbol", "zapfdingbats", "wingdings"])
    
    # Common bullet characters
    bullets = ["\u2022", "\u25E6", "\uf0b7", "\u2219", "\u2023", "\u2043"]
    if is_bullet_font or text[0] in bullets:
        return True, "bullet"
        
    # Numbered list pattern: 1. or 1) or (1)
    if re.match(r"^(\d+|[a-zA-Z]|[ivxIVX]+)[.)]\s", text) or re.match(r"^\(\d+\)\s", text):
        return True, "number"
        
    return False, None


@contextlib.contextmanager
def _temp_png(data: bytes):
    """Write *data* to a uniquely-named temp PNG, yield its path, then delete it."""
    fd, path = tempfile.mkstemp(suffix=".png")
    try:
        os.write(fd, data)
        os.close(fd)
        yield path
    finally:
        try:
            os.unlink(path)
        except OSError:
            pass


def _safe_pixmap_png(doc, xref: int) -> bytes | None:
    """Extract an image by xref and return PNG bytes, or None on failure.

    Handles CMYK and other colour spaces that fitz cannot encode directly as PNG
    by converting through Pillow instead.
    """
    try:
        pix = fitz.Pixmap(doc, xref)
        if pix.n - pix.alpha >= 4:
            # CMYK or other non-RGB/RGBA → convert to RGB via Pillow
            pil_img = Image.frombytes("CMYK", (pix.width, pix.height), pix.samples)
            pil_img = pil_img.convert("RGB")
            buf = io.BytesIO()
            pil_img.save(buf, format="PNG")
            return buf.getvalue()
        return pix.tobytes("png")
    except Exception as e:
        _log(f"  Pixmap extraction failed for xref {xref}: {e}")
        return None


def _build_docx_from_blocks(blocks: list[DocumentBlock], output_path: str) -> ConversionSummary:
    """Assemble a python-docx Document from extracted DocumentBlocks."""
    doc = Document()
    summary = ConversionSummary()
    
    current_page = 0
    for block in blocks:
        # Add page breaks if needed (handle skipped empty pages - Low #10)
        while block.page_num > current_page:
            doc.add_page_break()
            current_page += 1
            
        if block.block_type == "heading":
            level = min(block.level or 1, 9)
            doc.add_heading(block.text, level=level)
            summary.headings += 1
            
        elif block.block_type == "table":
            if not block.table_data:
                continue
            rows = len(block.table_data)
            cols = len(block.table_data[0]) if rows > 0 else 0
            if rows > 0 and cols > 0:
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Table Grid'
                for r_idx, row_data in enumerate(block.table_data):
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx < cols:
                            table.cell(r_idx, c_idx).text = str(cell_text or "")
                summary.tables += 1
                
        elif block.block_type == "image" or block.block_type == "scanned_page":
            if block.image_bytes:
                try:
                    img_stream = BytesIO(block.image_bytes)
                    # For scanned pages, we might want to make it full width
                    width = Inches(6.5) if block.block_type == "scanned_page" else None
                    doc.add_picture(img_stream, width=width)
                    if block.block_type == "scanned_page":
                        summary.scanned_pages += 1
                    else:
                        summary.images += 1
                except Exception as e:
                    _log(f"Error adding image to docx: {e}")
                    summary.skipped_elements.append(f"Image on page {block.page_num+1}")
                    
        elif block.block_type == "list_item":
            style = 'List Bullet' if block.list_style == "bullet" else 'List Number'
            p = doc.add_paragraph(style=style)
            for span in block.spans:
                run = p.add_run(span.text)
                run.bold = span.bold
                run.italic = span.italic
                run.font.size = Pt(span.font_size)
            summary.list_items += 1
            
        else: # paragraph
            p = doc.add_paragraph()
            # Basic alignment heuristic based on bbox (FR-005)
            # This is a bit crude in docx without knowing page width, but we'll stick to styles
            for span in block.spans:
                run = p.add_run(span.text)
                run.bold = span.bold
                run.italic = span.italic
                run.font.size = Pt(span.font_size)
            summary.text_blocks += 1

    doc.save(output_path)
    return summary


def _md_escape(text: str) -> str:
    """Escape the characters that would otherwise be read as markup."""
    out = text
    for ch in ("\\", "`", "*", "_", "[", "]", "#"):
        out = out.replace(ch, "\\" + ch)
    return out


def _md_spans(spans: list[SpanInfo], fallback: str | None) -> str:
    """Spans as inline markdown."""
    if not spans:
        return _md_escape(fallback or "")
    parts: list[str] = []
    for span in spans:
        text = span.text
        if not text:
            continue
        # Emphasis markers cannot touch whitespace in markdown, or the
        # asterisks render literally. Keep padding outside the marks.
        core = text.strip()
        if not core:
            parts.append(text)
            continue
        lead = text[:len(text) - len(text.lstrip())]
        trail = text[len(text.rstrip()):]
        marks = ("**" if span.bold else "") + ("*" if span.italic else "")
        parts.append(lead + marks + _md_escape(core) + marks[::-1] + trail)
    return "".join(parts)


def _md_table(rows: list[list[str]]) -> str:
    """A GFM pipe table, treating the first row as the header."""
    width = max((len(r) for r in rows), default=0)
    if not width:
        return ""

    def line(cells) -> str:
        padded = [str(c or "").replace("|", "\\|").replace("\n", " ")
                  for c in list(cells) + [""] * (width - len(cells))]
        return "| " + " | ".join(padded) + " |"

    out = [line(rows[0]), "|" + "|".join([" --- "] * width) + "|"]
    out.extend(line(r) for r in rows[1:])
    return "\n".join(out)


_LIST_MARKER_RE = re.compile(
    r"^\s*(?:[\u2022\u25CF\u25AA\u2023\u2043o]|\\?[-*+]|\(?\d{1,3}[.)]|"
    r"\(?[a-zA-Z][.)])\s+")


def _strip_list_marker(text: str) -> str:
    """Remove a bullet/number the source already carried.

    The PDF extractor keeps the literal glyph in the block text, so emitting our
    own marker on top gives "- \u2022 first" or "1. 1. one".
    """
    return _LIST_MARKER_RE.sub("", text, count=1).strip()


def _build_md_from_blocks(blocks: list[DocumentBlock], output_path: str) -> ConversionSummary:
    """Serialise DocumentBlocks to Markdown.

    Reuses the blocks the DOCX path already consumes, so heading/list/table
    detection is shared rather than reimplemented.

    Images are written to a `<stem>_images/` folder beside the .md and linked
    relatively, since Markdown cannot embed bytes. Scanned pages still get a
    warning so a near-empty result is explained.
    """
    summary = ConversionSummary()
    lines: list[str] = []
    stem = os.path.splitext(os.path.basename(output_path))[0]
    img_dir_name = "%s_images" % stem
    img_dir: str | None = None
    img_count = 0
    # Markdown needs a blank line between blocks of different kinds, or they
    # merge: a table right after a list item becomes part of that item, and a
    # numbered list touching a bulleted one joins it. Tracked rather than
    # emitted eagerly so consecutive list items stay tight.
    prev_kind: str | None = None

    def gap(kind: str) -> None:
        nonlocal prev_kind
        if prev_kind is not None and prev_kind != kind:
            if lines and lines[-1] != "":
                lines.append("")
            # A blank line alone does NOT split two adjacent lists - a bulleted
            # list followed by a numbered one renders as a single <ul>. An HTML
            # comment is the standard separator that forces a new list while
            # staying invisible when rendered.
            if prev_kind.startswith("list-") and kind.startswith("list-"):
                lines.extend(("<!-- -->", ""))
        prev_kind = kind

    for block in blocks:
        if block.block_type == "heading":
            gap("heading")
            level = min(max(block.level or 1, 1), 6)
            # No emphasis inside a heading: PDF heading spans are usually bold,
            # which would render as "## **Title**".
            text = _md_escape((block.text or "").strip()) or _md_spans(
                block.spans, block.text)
            lines += ["#" * level + " " + text.strip(), ""]
            summary.headings += 1

        elif block.block_type == "table":
            if block.table_data:
                gap("table")
                lines += [_md_table(block.table_data), ""]
                summary.tables += 1

        elif block.block_type in ("image", "scanned_page"):
            if block.block_type == "scanned_page":
                summary.scanned_pages += 1
                summary.warnings.append(
                    "Page %d is a scanned image - it has no text to extract."
                    % (block.page_num + 1))
            if not block.image_bytes:
                continue
            # Markdown cannot embed bytes, so pictures go in a sibling folder
            # and are linked relatively — the .md stays portable if the pair
            # is moved together.
            try:
                if img_dir is None:
                    img_dir = os.path.join(
                        os.path.dirname(output_path) or ".", img_dir_name)
                os.makedirs(img_dir, exist_ok=True)
                img_count += 1
                name = "img%03d.%s" % (img_count, block.image_ext or "png")
                with open(os.path.join(img_dir, name), "wb") as fh:
                    fh.write(block.image_bytes)
                gap("image")
                # Forward slashes: a markdown link is a URL, not a Windows path.
                lines += ["![page %d](%s/%s)" % (
                    block.page_num + 1, img_dir_name, name), ""]
                summary.images += 1
            except OSError as e:
                _log("Could not write markdown image: %s" % e)
                summary.skipped_elements.append(
                    "Image on page %d" % (block.page_num + 1))

        elif block.block_type == "list_item":
            style = "number" if block.list_style == "number" else "bullet"
            gap("list-" + style)
            indent = "  " * max((block.level or 1) - 1, 0)
            marker = "1." if style == "number" else "-"
            body = _strip_list_marker(_md_spans(block.spans, block.text))
            lines.append(indent + marker + " " + body)
            summary.list_items += 1

        else:
            text = _md_spans(block.spans, block.text)
            if text.strip():
                gap("paragraph")
                lines += [text, ""]
                summary.text_blocks += 1

    with open(output_path, "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines).rstrip() + "\n")
    return summary


def _pdf_to_md(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """PDF to Markdown - same extraction as PDF to DOCX, different serialiser."""
    try:
        with fitz.open(input_path) as doc:
            if doc.is_encrypted:
                return False, "PDF is encrypted/password-protected", None

            total_pages = len(doc)
            body_font_size = _detect_body_font_size(doc)
            all_blocks = []
            extracted_images = set()
            extraction_warnings = []

            for page_num, page in enumerate(doc):
                if cancel_event and cancel_event.is_set():
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    return False, "Conversion cancelled by user", None
                if progress_callback:
                    progress_callback(page_num + 1, total_pages)
                all_blocks.extend(_extract_page_blocks(
                    page, page_num, body_font_size, doc, extracted_images,
                    warnings=extraction_warnings))

            summary = _build_md_from_blocks(all_blocks, output_path)
            summary.total_pages = total_pages
            if extraction_warnings:
                summary.warnings.extend(extraction_warnings)
            return True, output_path, summary
    except Exception as e:
        _log("PDF to Markdown failed: %s" % e)
        return False, str(e), None


def _docx_to_blocks(input_path: str, cancel_event=None) -> list[DocumentBlock] | None:
    """Read a DOCX into DocumentBlocks. Returns None if cancelled.

    Word records heading level, list style and tables explicitly, so nothing
    needs inferring here - unlike PDF, where layout has to be guessed.
    """
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(input_path)
    blocks: list[DocumentBlock] = []

    # Walk the body in document order so tables stay between the correct
    # paragraphs; doc.paragraphs and doc.tables are separate flat lists.
    for child in doc.element.body.iterchildren():
        if cancel_event and cancel_event.is_set():
            return None
        tag = child.tag.split("}")[-1]

        if tag == "p":
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            spans = [SpanInfo(text=r.text, bold=bool(r.bold),
                              italic=bool(r.italic), font_size=0.0,
                              font_name="", color=(0, 0, 0))
                     for r in para.runs if r.text]
            style = (para.style.name or "") if para.style is not None else ""
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 1
                blocks.append(DocumentBlock("heading", text=text,
                                            spans=spans, level=level))
            elif style.startswith("List"):
                blocks.append(DocumentBlock(
                    "list_item", text=text, spans=spans, level=1,
                    list_style="number" if "Number" in style else "bullet"))
            else:
                blocks.append(DocumentBlock("paragraph", text=text,
                                            spans=spans))

        elif tag == "tbl":
            rows = [[c.text.strip() for c in r.cells]
                    for r in Table(child, doc).rows]
            if rows:
                blocks.append(DocumentBlock("table", table_data=rows))

    return blocks


def _docx_to_md(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """DOCX to Markdown."""
    try:
        blocks = _docx_to_blocks(input_path, cancel_event)
        if blocks is None:
            return False, "Conversion cancelled by user", None
        if progress_callback:
            progress_callback(1, 1)
        summary = _build_md_from_blocks(blocks, output_path)
        summary.total_pages = 1
        return True, output_path, summary
    except Exception as e:
        _log("DOCX to Markdown failed: %s" % e)
        return False, str(e), None


# Slide budget. A 10 x 7.5in slide with the default 18pt body font fits roughly
# this much before text spills off the bottom, which is what makes a naive
# "one PDF page = one slide" conversion look broken. Blocks are packed up to
# these limits and then continued on a new slide.
_SLIDE_MAX_LINES = 11          # body lines per slide
_SLIDE_MAX_CHARS = 620         # characters per slide, for long paragraphs
_SLIDE_CHARS_PER_LINE = 58     # to estimate how many lines a paragraph wraps to
_PPTX_LAYOUT_TITLE_CONTENT = 1
_PPTX_LAYOUT_TITLE_ONLY = 5
_PPTX_LAYOUT_BLANK = 6
_PPTX_TABLE_MAX_ROWS = 12      # taller tables are split across slides


def _estimate_lines(text: str) -> int:
    """How many rendered lines *text* will wrap to in the body placeholder."""
    return max(1, -(-len(text) // _SLIDE_CHARS_PER_LINE))


# A "design deck" PDF is already slide-shaped: every page is a composed layout
# with a full-bleed background image and text set on top of it. Decomposing
# such a page into blocks and re-flowing them destroys the design — the
# background lands on one slide, the overlaid title is orphaned, and the result
# looks broken. Those pages are rendered whole instead.
_DECK_MIN_ASPECT = 1.4         # 4:3 is 1.33; anything wider is presentation-shaped
_DECK_BG_COVERAGE = 0.85       # an image covering this much of the page is a background
# 150 DPI on a 1920x1080 page gives a 4000 px PNG — 36 of those made a 394 MB
# deck. A slide is displayed at ~1920 px wide at most, so cap the long edge
# there and encode photographic pages as JPEG: same visual result, ~2% the size.
# A drop shadow is a small image lying on top of the words; the background
# photo is a big image that barely touches them.
_SHADOW_MAX_PAGE_FRAC = 0.25      # of the page area
_SHADOW_MIN_TEXT_OVERLAP = 0.5    # of the image's own area

# Render the page ABOVE its nominal point size. A 1920 pt page rendered at
# zoom 1.0 is 1920 px — i.e. 72 DPI — and the deck's own photos are placed at
# ~1.06 native pixels per point, so 1:1 was throwing away source detail and
# JPEG-85 on top of that was visibly soft. 2880 px with quality 92 sits
# comfortably above the source resolution; measured on the 36-page sample
# deck that is ~25 MB against a 44 MB original.
_DECK_MAX_PX = 2880
_DECK_JPEG_QUALITY = 92


def _page_is_designed_slide(page) -> bool:
    """True when *page* is a composed slide rather than a flowed document page.

    Two signals together: the page is presentation-shaped (wide), and it carries
    an image large enough to be a background. Text over a full-bleed image is
    exactly the layout that block extraction cannot reassemble.
    """
    rect = page.rect
    area = rect.get_area()
    if area <= 0 or rect.height <= 0:
        return False
    if rect.width / rect.height < _DECK_MIN_ASPECT:
        return False
    try:
        images = page.get_images(full=True)
    except Exception:
        return False
    for img in images:
        bbox = _image_bbox(page, img)
        if bbox is None:
            continue
        # Clip to the page: these backgrounds are deliberately oversized
        # (measured coverage up to 308% before clipping).
        visible = fitz.Rect(bbox) & rect
        if visible.is_empty:
            continue
        if visible.get_area() / area >= _DECK_BG_COVERAGE:
            return True
    return False


_PARA_BREAK_EM = 0.45     # em; gap above this is a paragraph break, not line spacing
_TRACK_MIN_EM = 0.045     # em; below this PowerPoint's default is close enough


def _fix_rtl_order(text: str) -> str:
    """Turn a visually-ordered Arabic/Hebrew run back into logical order.

    PDF stores RTL text the way it is painted — right to left, using Arabic
    *presentation forms* (U+FE70..U+FEFF) rather than base letters. Extraction
    hands that back verbatim, so the string is reversed and unshaped, and
    pasting it into PowerPoint gives mirrored gibberish.

    NFKC maps each presentation form to its base letter; reversing the
    characters within each word AND the order of the words undoes the visual
    layout. Verified round-trip: the visual form of "مرحبا بكم في زنجبار"
    normalises and reverses back to exactly that string.

    A dedicated bidi library would handle mixed-direction lines more
    rigorously, but this needs no new dependency and is correct for the
    single-direction headings and paragraphs decks actually contain.
    """
    import unicodedata

    if not _is_rtl(text):
        return text
    # NBSP is what PDF word gaps usually come back as.
    normalised = unicodedata.normalize("NFKC", text).replace(" ", " ")
    lead = " " if normalised[:1] == " " else ""
    tail = " " if normalised[-1:] == " " else ""
    words = [w for w in normalised.split(" ") if w]
    return lead + " ".join(w[::-1] for w in reversed(words)) + tail


def _measure_tracking(chars, size: float) -> tuple[str, float]:
    """Recover a span's real wording and its letter spacing, in em.

    Design decks spell out tracking in two different ways, and passing either
    through verbatim gives the wrong result in PowerPoint:

    1. **Fake tracking** — a literal space between every letter, two between
       words: ``"T O L E O  T O W E R"``. PowerPoint renders those spaces at
       its own width, so word gaps collapse and the line reads as one blob.
       The wording has to be rebuilt ("TOLEO TOWER") and the spacing re-applied
       as real tracking, or the words run together.
    2. **Real tracking** — no space characters at all, spacing carried in the
       glyph advances. Extraction reports a solid run ("TOLEOTOWER"), which
       PowerPoint then sets tight.

    Returns ``(text, em)`` where *em* is 0.0 when the span is ordinary text.
    """
    import statistics

    text = "".join(c["c"] for c in chars)
    letters = [c for c in chars if c["c"] != " "]
    if size <= 0 or len(letters) < 3:
        return text, 0.0

    gaps = [b["bbox"][0] - a["bbox"][2] for a, b in zip(letters, letters[1:])]
    gaps = [g for g in gaps if g > -size]      # ignore a line/column wrap
    if not gaps:
        return text, 0.0
    em = statistics.median(gaps) / size

    # Fake tracking: most of the characters are spaces. Collapse runs of two
    # or more spaces to one word break and drop the single separators.
    stripped = text.replace(" ", "")
    if stripped and (len(text) - len(stripped)) >= len(stripped) * 0.6:
        words, current = [], []
        run = 0
        for ch in text:
            if ch == " ":
                run += 1
                continue
            if run >= 2 and current:
                words.append("".join(current))
                current = []
            run = 0
            current.append(ch)
        if current:
            words.append("".join(current))
        rebuilt = " ".join(w for w in words if w)
        lead = " " if text[:1] == " " else ""
        tail = " " if text[-1:] == " " else ""
        # The gap that matters is letter edge to letter edge INCLUDING the
        # separator glyph the designer typed. Measuring only the space either
        # side of it understates the tracking (0.166 em measured where the
        # design actually sets 0.28 em).
        adv = [b["bbox"][0] - a["bbox"][2]
               for a, b in zip(letters, letters[1:])
               if 0 < b["bbox"][0] - a["bbox"][2] < size]
        if adv:
            em = statistics.median(adv) / size
        return lead + rebuilt + tail, max(0.0, em)

    return text, em if em >= _TRACK_MIN_EM else 0.0


def _page_text_overlays(page) -> list[dict]:
    """Every text LINE on *page* as a positioned, styled overlay spec.

    One spec per line, not per block. A PDF text block can contain lines of
    different sizes at different x offsets — page 24 of the sample deck holds
    58 pt "RESIDENTIAL", 47 pt "9" and 19.7/33.8 pt "th FLOOR" in one block —
    and pouring those into a single textbox re-wraps them into overlapping
    nonsense. Placing each line where the PDF put it reproduces the design.

    Consecutive lines that share a size and left edge are merged back into one
    box so a genuine paragraph stays a single editable object.
    """
    rect = page.rect
    if rect.width <= 0 or rect.height <= 0:
        return []
    try:
        # rawdict, not dict: per-character boxes are needed to recover the
        # design's letter spacing (see _measure_tracking).
        data = page.get_text("rawdict")
    except Exception:
        return []

    raw: list[dict] = []
    for block in data.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            runs = []
            lx0 = ly0 = lx1 = ly1 = None
            for span in line.get("spans", []):
                size = float(span.get("size", 12.0))
                text, track = _measure_tracking(span.get("chars", []), size)
                if not text.strip():
                    continue
                text = _fix_rtl_order(text)
                sx0, sy0, sx1, sy1 = span["bbox"]
                lx0 = sx0 if lx0 is None else min(lx0, sx0)
                ly0 = sy0 if ly0 is None else min(ly0, sy0)
                lx1 = sx1 if lx1 is None else max(lx1, sx1)
                ly1 = sy1 if ly1 is None else max(ly1, sy1)
                font = span.get("font") or ""
                runs.append({
                    "text": text,
                    "size": size,
                    "track": track,
                    "color": int(span.get("color", 0)) & 0xFFFFFF,
                    "bold": "bold" in font.lower() or "black" in font.lower(),
                    "italic": "italic" in font.lower() or "oblique" in font.lower(),
                })
            if runs:
                raw.append({"runs": runs, "x0": lx0, "y0": ly0,
                            "x1": lx1, "y1": ly1,
                            "size": max(r["size"] for r in runs)})

    # Merge consecutive lines that read as one paragraph, so a real paragraph
    # stays a single editable object. Three conditions, each learned from the
    # sample deck:
    #   * same size — a 58 pt heading and a 34 pt subtitle are separate things
    #   * aligned — sharing a left edge OR a midpoint (centred text shares
    #     neither left nor right, only its centre)
    #   * stacked at line spacing — the gap may be slightly NEGATIVE, because
    #     glyph bboxes include ascenders/descenders that overlap between lines
    #     (measured -5.5 pt on a 57.9 pt heading), so a `gap >= 0` test wrongly
    #     splits a two-line heading.
    groups: list[list[dict]] = []
    for line in raw:
        prev = groups[-1][-1] if groups else None
        same_para = False
        if prev is not None and abs(line["size"] - prev["size"]) < 0.6:
            aligned = (
                abs(line["x0"] - prev["x0"]) < 2.0
                or abs(((line["x0"] + line["x1"]) / 2)
                       - ((prev["x0"] + prev["x1"]) / 2)) < 2.0
            )
            gap = line["y0"] - prev["y1"]
            stacked = -line["size"] * 0.35 <= gap <= line["size"] * 1.2
            same_para = aligned and stacked
        if same_para:
            # Record how far this line sits below the previous one. Lines of
            # one paragraph are tight (measured -0.01 em, glyph boxes touch);
            # a paragraph BREAK inside the same run of text shows up as a
            # clear jump (0.89 em on the sample deck). Keeping the run merged
            # preserves it as one editable object while still reproducing the
            # blank line the design has.
            line["gap_em"] = (line["y0"] - prev["y1"]) / max(1.0, line["size"])
            groups[-1].append(line)
        else:
            line["gap_em"] = 0.0
            groups.append([line])

    out: list[dict] = []
    for group in groups:
        x0 = min(l["x0"] for l in group)
        y0 = min(l["y0"] for l in group)
        x1 = max(l["x1"] for l in group)
        y1 = max(l["y1"] for l in group)

        align = "left"
        if len(group) > 1:
            mids = [(l["x0"] + l["x1"]) / 2 for l in group]
            starts = [l["x0"] for l in group]
            if max(starts) - min(starts) > 2.0 and \
                    (max(mids) - min(mids)) < (max(starts) - min(starts)):
                align = "center"
        else:
            page_mid = rect.width / 2
            line_mid = (x0 + x1) / 2
            if abs(line_mid - page_mid) < rect.width * 0.02 \
                    and (x1 - x0) < rect.width * 0.9:
                align = "center"

        out.append({
            "align": align,
            "x": x0 / rect.width,
            "y": y0 / rect.height,
            "w": (x1 - x0) / rect.width,
            "h": (y1 - y0) / rect.height,
            "page_h": rect.height,
            "lines": [{"runs": l["runs"],
                       "space_before": l.get("gap_em", 0.0) > _PARA_BREAK_EM}
                      for l in group],
        })
    return out


def _render_page_background(page, page_num: int):
    """Strip *page*'s headline text so the editable overlays are not doubled.

    Design work carries "text" in three forms and each needs its own removal:

    * **Live text** — taken by the redaction itself.
    * **Text converted to CURVES** — this deck draws several headlines as one
      stroke-only path per glyph, and ``apply_redactions`` leaves vector art
      alone, so those outlines showed through the overlay as hollow letters.
      They are erased individually by ``_outlined_glyph_paths``.
    * **A drop SHADOW** — a small soft-masked image lying on the words.
      Redaction leaves it behind, and being a blurred copy of the same
      lettering it reads as grey shading under the overlay.

    Removal is deliberately targeted. ``apply_redactions(graphics=2)`` looks
    like the tidy answer but destroys the page: it drops any art *overlapping*
    a redaction rectangle, so a full-page colour field that clips a text block
    by 4%, or a palette swatch that touches a caption, disappears with it — a
    colour-system slide came out blank white.
    """
    import fitz as _fitz

    text_rects: list = []
    try:
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            box = _fitz.Rect(block["bbox"])
            if not box.is_empty:
                text_rects.append(box)
    except Exception as exc:
        _log("text scan failed on page %d: %s" % (page_num + 1, exc))

    # Outlined lettering has to be covered rather than deleted: PyMuPDF cannot
    # erase a vector path, and a redaction rectangle only whites the area out
    # before the path is drawn over it again.
    #
    # A patch is only honest where the background is flat. Over a photograph
    # every option looked worse than leaving the design alone: a flat sampled
    # colour shows as a rectangle per letter, and copying a slice of nearby
    # background stamps in unrelated scenery. So a block whose outlined
    # lettering cannot be covered cleanly keeps its artwork, and the caller
    # drops that block's overlay so the words are not printed twice.
    # Decide per BLOCK, never per glyph: covering some letters of a heading and
    # not others leaves a half-erased headline with the overlay on top, which
    # is worse than either extreme. Sample every glyph of a block first, and
    # only patch if all of them can be hidden.
    by_block: dict = {}
    for rect, block in _outlined_glyph_paths(page, text_rects):
        by_block.setdefault(tuple(block), []).append(rect)

    unpatchable: list = []
    for key, rects in by_block.items():
        colours = []
        for rect in rects:
            try:
                colours.append(_sample_ring_colour(page, rect))
            except Exception as exc:
                _log("glyph sampling failed on page %d: %s"
                     % (page_num + 1, exc))
                colours.append(None)
        if any(c is None for c in colours):
            unpatchable.append(_fitz.Rect(key))
            continue
        for rect, colour in zip(rects, colours):
            try:
                page.draw_rect(rect, color=colour, fill=colour, width=0)
            except Exception as exc:
                _log("glyph cover failed on page %d: %s"
                     % (page_num + 1, exc))

    text_rects = [r for r in text_rects
                  if not any(r == u for u in unpatchable)]

    _remove_text_shadows(page, page_num, text_rects)

    for box in text_rects:
        page.add_redact_annot(box)

    try:
        # IMAGE_NONE keeps the background photo; graphics=0 leaves every vector
        # shape alone, because the page's colour fields and swatches are vector
        # art and graphics=2 wipes them (see the docstring).
        page.apply_redactions(images=_fitz.PDF_REDACT_IMAGE_NONE, graphics=0)
    except TypeError:
        try:
            page.apply_redactions(images=_fitz.PDF_REDACT_IMAGE_NONE)
        except Exception as exc:
            _log("redaction failed on page %d: %s" % (page_num + 1, exc))
    except Exception as exc:
        _log("redaction failed on page %d: %s" % (page_num + 1, exc))
    return unpatchable


# An outlined glyph is a stroke-only path, about text height, sitting almost
# entirely inside a reported text block. A design swatch is filled, hundreds of
# points tall, and only clips a text block's edge.
_PATCH_MAX_SPREAD = 40      # channel range a flat patch may hide in
_GLYPH_MIN_CONTAINMENT = 0.9
_GLYPH_MAX_HEIGHT_RATIO = 1.6     # of the enclosing text block's height


def _outlined_glyph_paths(page, text_rects) -> list:
    """``(rect, text_block)`` for each vector path that is a letter of a heading.

    Each rectangle is grown by the path's stroke width: ``get_drawings``
    reports the path's own bbox, and a stroke of width *w* paints w/2 beyond it
    on every side, so a patch cut to the reported rect leaves a thin white
    outline round itself — which looked like a box drawn round each letter.
    """
    import fitz as _fitz

    out = []
    if not text_rects:
        return out
    try:
        drawings = page.get_drawings()
    except Exception:
        return out
    for drawing in drawings:
        rect = _fitz.Rect(drawing["rect"])
        if rect.is_empty or rect.is_infinite or rect.height <= 0:
            continue
        area = rect.get_area()
        if area <= 0:
            continue
        # Filled shapes are design elements; letters here are stroked outlines.
        if drawing.get("type") == "f":
            continue
        for block in text_rects:
            if (rect & block).get_area() <= area * _GLYPH_MIN_CONTAINMENT:
                continue
            if rect.height > block.height * _GLYPH_MAX_HEIGHT_RATIO:
                continue
            grow = max(1.0, float(drawing.get("width") or 0.0))
            out.append((rect + (-grow, -grow, grow, grow), block))
            break
    return out


def _remove_text_shadows(page, page_num: int, text_rects) -> None:
    """Delete a headline's drop-shadow image layer.

    A shadow is small and lies mostly inside a text block; the background
    photograph covers the whole page and barely touches the text (measured
    0.05 of the page with 0.61 overlap versus 1.46 with 0.02).
    """
    if not text_rects:
        return
    page_area = page.rect.get_area() or 1.0
    try:
        images = page.get_images(full=True)
    except Exception:
        return
    for img in images:
        try:
            box = _image_bbox(page, img)
        except Exception:
            continue
        if box is None or box.is_empty or box.is_infinite:
            continue
        area = box.get_area()
        if area <= 0 or area > page_area * _SHADOW_MAX_PAGE_FRAC:
            continue
        inside = max(((box & t).get_area() / area for t in text_rects),
                     default=0.0)
        if inside > _SHADOW_MIN_TEXT_OVERLAP:
            try:
                page.delete_image(img[0])
            except Exception as exc:
                _log("shadow removal failed on page %d: %s"
                     % (page_num + 1, exc))


def _sample_ring_colour(page, box):
    """Median colour of a thin ring around *box*, as (r, g, b) in 0-1.

    Used to cover an outlined glyph. Returns None when the surroundings are
    too varied for a flat patch to be honest, in which case the caller leaves
    the glyph alone rather than stamping a visible rectangle onto artwork.
    """
    import fitz as _fitz

    rect = page.rect
    margin = max(2.0, box.height * 0.25)
    outer = (box + (-margin, -margin, margin, margin)) & rect
    if outer.is_empty:
        return None
    try:
        pix = page.get_pixmap(clip=outer, alpha=False)
    except Exception:
        return None
    if pix.width < 3 or pix.height < 3 or pix.n < 3:
        return None

    samples = []
    for x in range(pix.width):
        for y in (0, pix.height - 1):
            samples.append(pix.pixel(x, y))
    for y in range(pix.height):
        for x in (0, pix.width - 1):
            samples.append(pix.pixel(x, y))
    if not samples:
        return None
    chans = [sorted(sm[i] for sm in samples) for i in range(3)]
    mid = len(samples) // 2
    spread = max(c[-1] - c[0] for c in chans)
    # 40, not 90: at 90 a soft photographic gradient still qualified and each
    # letter's patch showed as a faint rectangle. Anything above this leaves
    # the artwork alone (and the caller drops that block's overlay).
    if spread > _PATCH_MAX_SPREAD:
        return None
    return tuple(c[mid] / 255.0 for c in chans)


def _overlay_in_blocks(spec, page_rect, blocks) -> bool:
    """True when *spec* covers one of *blocks* (page fractions vs points)."""
    import fitz as _fitz

    box = _fitz.Rect(spec["x"] * page_rect.width,
                     spec["y"] * page_rect.height,
                     (spec["x"] + spec["w"]) * page_rect.width,
                     (spec["y"] + spec["h"]) * page_rect.height)
    area = box.get_area()
    if area <= 0:
        return False
    return any((box & b).get_area() > area * 0.5 for b in blocks)


def _render_page_block(page, page_num: int, editable_text: bool = True) -> DocumentBlock | None:
    """The whole page as one image block, preserving its composed layout.

    editable_text: strip the text from the render and return it as overlay
    specs, so the slide keeps its design AND its text stays editable.
    """
    rect = page.rect
    longest = max(rect.width, rect.height) or 1
    # Scale to the display budget rather than a fixed DPI: a 1920 pt page and
    # an A4 page then produce similarly sized images.
    zoom = min(_DECK_MAX_PX / longest, 4.0)

    overlays: list[dict] = []
    if editable_text:
        overlays = _page_text_overlays(page)
        if overlays:
            kept = _render_page_background(page, page_num)
            if kept:
                # These blocks still show their original lettering, because it
                # could not be removed without damaging the artwork. Drop
                # their overlays so the words are not printed twice.
                overlays = [o for o in overlays
                            if not _overlay_in_blocks(o, page.rect, kept)]

    try:
        import fitz as _fitz
        pix = page.get_pixmap(matrix=_fitz.Matrix(zoom, zoom), alpha=False)
        data = pix.tobytes("jpeg", jpg_quality=_DECK_JPEG_QUALITY)
        ext = "jpg"
    except Exception as exc:
        _log("page render failed on page %d: %s" % (page_num + 1, exc))
        return None
    if not data:
        return None
    block = DocumentBlock(
        block_type="slide_image",
        image_bytes=data,
        image_ext=ext,
        bbox=tuple(rect),
        page_num=page_num,
    )
    # Stashed on the block so the serialiser needs no extra plumbing.
    block.spans = []
    block.table_data = None
    block.text_overlays = overlays
    return block


def _group_blocks_into_slides(blocks: list[DocumentBlock]) -> list[dict]:
    """Pack DocumentBlocks into slide-sized groups.

    A heading starts a new slide and becomes its title; body blocks accumulate
    until the slide is full, then continue on another slide carrying the same
    title. Tables get their own slide because they need the full canvas.

    Returns [{"title": str|None, "blocks": [...], "continued": bool}].
    """
    slides: list[dict] = []
    title: str | None = None
    current: list[DocumentBlock] = []
    lines = 0
    chars = 0
    continued = False

    def flush(more: bool = False) -> None:
        nonlocal current, lines, chars, continued
        if current or title is not None:
            slides.append({"title": title, "blocks": current,
                           "continued": continued, "caption": None})
        current = []
        lines = 0
        chars = 0
        continued = more

    for block in blocks:
        if block.block_type == "heading":
            flush()
            title = (block.text or "").strip() or None
            continued = False
            continue

        if block.block_type == "slide_image":
            # A composed page is the slide. No title, no caption, nothing else
            # on it — anything added would sit on top of the artwork.
            flush()
            slides.append({"title": None, "blocks": [block],
                           "continued": False, "caption": None,
                           "full_bleed": True})
            continue

        if block.block_type in ("image", "scanned_page"):
            # A picture needs the slide canvas, so it gets its own slide — the
            # same treatment tables get, and for the same reason.
            if not block.image_bytes:
                continue
            was_first = not slides and not current
            flush()
            slides.append({"title": title, "blocks": [block],
                           "continued": not was_first, "caption": None})
            continued = True
            continue

        if block.block_type == "table":
            rows = block.table_data or []
            if not rows:
                continue
            flush()
            # Split a tall table so it never runs off the slide, repeating the
            # header row on each part.
            header, body = rows[0], rows[1:]
            if len(rows) <= _PPTX_TABLE_MAX_ROWS:
                chunks = [rows]
            else:
                size = _PPTX_TABLE_MAX_ROWS - 1
                chunks = [[header] + body[i:i + size]
                          for i in range(0, len(body), size)]
            for i, chunk in enumerate(chunks):
                slides.append({
                    "title": title,
                    "blocks": [DocumentBlock("table", table_data=chunk,
                                             page_num=block.page_num)],
                    "continued": i > 0,
                    "caption": None,
                })
            continued = True   # anything after the table is a continuation
            continue

        text = (block.text or "").strip()
        if not text:
            continue
        need = _estimate_lines(text)
        if current and (lines + need > _SLIDE_MAX_LINES
                        or chars + len(text) > _SLIDE_MAX_CHARS):
            flush(more=True)
        current.append(block)
        lines += need
        chars += len(text)

    flush()

    # A short paragraph immediately after a figure is its caption; move it onto
    # the figure's own slide rather than leaving it stranded.
    folded: list[dict] = []
    for group in slides:
        prev = folded[-1] if folded else None
        is_caption_candidate = (
            prev is not None
            and any(b.block_type in ("image", "scanned_page")
                    for b in prev["blocks"])
            and prev.get("caption") is None
            and len(group["blocks"]) == 1
            and group["blocks"][0].block_type == "paragraph"
            and len((group["blocks"][0].text or "")) <= 160
        )
        if is_caption_candidate:
            prev["caption"] = (group["blocks"][0].text or "").strip()
            continue
        folded.append(group)
    slides = folded

    # A trailing flush can leave a title-only continuation slide with no body
    # (e.g. the last real content was a table). Drop any slide that carries no
    # blocks unless it is a genuine section header — i.e. the FIRST slide under
    # that title, not a "(cont.)" of it.
    return [g for g in slides
            if g["blocks"] or g.get("caption")
            or (g["title"] and not g["continued"])]


def _add_table_slide(prs, slide, rows: list[list[str]]) -> None:
    """Place a table on *slide*, sized to the slide and clamped to sane widths."""
    from pptx.util import Emu, Pt

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    left = Emu(int(prs.slide_width * 0.06))
    top = Emu(int(prs.slide_height * 0.28))
    width = Emu(int(prs.slide_width * 0.88))
    # Let PowerPoint grow rows if text wraps; this is a starting height only.
    height = Emu(int(min(prs.slide_height * 0.62, Pt(24).emu * n_rows)))

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    for r, row in enumerate(rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)


def _add_picture_slide(prs, slide, block: DocumentBlock) -> bool:
    """Place a picture on *slide*, fitted to the content area, centred.

    Scales by the tighter of the two ratios so the image never stretches and
    never spills off the slide.
    """
    from io import BytesIO

    from PIL import Image as _PILImage
    from pptx.util import Emu

    avail_w = int(prs.slide_width * 0.88)
    avail_h = int(prs.slide_height * 0.62)
    top_margin = int(prs.slide_height * 0.28)

    try:
        with _PILImage.open(BytesIO(block.image_bytes)) as im:
            px_w, px_h = im.size
    except Exception:
        px_w = px_h = 0
    if px_w <= 0 or px_h <= 0:
        return False

    scale = min(avail_w / px_w, avail_h / px_h)
    width = int(px_w * scale)
    height = int(px_h * scale)
    left = int((prs.slide_width - width) / 2)
    top = int(top_margin + (avail_h - height) / 2)

    try:
        slide.shapes.add_picture(BytesIO(block.image_bytes), Emu(left),
                                 Emu(top), Emu(width), Emu(height))
        return True
    except Exception as exc:
        _log("Could not place picture: %s" % exc)
        return False


def _match_slide_size(prs, blocks: list[DocumentBlock]) -> None:
    """Adopt the source page shape for a deck, so nothing has to be cropped."""
    from pptx.util import Emu

    page = next((b for b in blocks if b.block_type == "slide_image"), None)
    if page is None:
        return
    x0, y0, x1, y1 = page.bbox
    width, height = x1 - x0, y1 - y0
    if width <= 0 or height <= 0:
        return
    # Keep the template's width and derive the height, so font sizes elsewhere
    # stay sensible relative to the canvas.
    prs.slide_height = Emu(int(prs.slide_width * height / width))


def _add_full_bleed_picture(prs, slide, block: DocumentBlock) -> bool:
    """Cover the whole slide with *block*'s image, centre-cropping any overflow.

    Uses the LARGER scale factor (cover, not contain) so there are no letterbox
    bars; a 16:9 page on a 16:9 slide is an exact fit either way.
    """
    from io import BytesIO

    from PIL import Image as _PILImage
    from pptx.util import Emu

    try:
        with _PILImage.open(BytesIO(block.image_bytes)) as im:
            px_w, px_h = im.size
    except Exception:
        return False
    if px_w <= 0 or px_h <= 0:
        return False

    scale = max(prs.slide_width / px_w, prs.slide_height / px_h)
    width = int(px_w * scale)
    height = int(px_h * scale)
    left = int((prs.slide_width - width) / 2)
    top = int((prs.slide_height - height) / 2)
    try:
        slide.shapes.add_picture(BytesIO(block.image_bytes), Emu(left),
                                 Emu(top), Emu(width), Emu(height))
        return True
    except Exception as exc:
        _log("Could not place full-bleed picture: %s" % exc)
        return False


_OVERLAY_FONT = "Calibri"


def _is_rtl(text: str) -> bool:
    """True when *text* contains Arabic, Hebrew, Syriac or Thaana letters."""
    return any(
        "֐" <= ch <= "׿"      # Hebrew
        or "؀" <= ch <= "ۿ"   # Arabic
        or "܀" <= ch <= "ݏ"   # Syriac
        or "ހ" <= ch <= "޿"   # Thaana
        or "ࢠ" <= ch <= "ࣿ"   # Arabic Extended-A
        or "ﭐ" <= ch <= "﷿"   # Arabic Presentation Forms-A
        or "ﹰ" <= ch <= "﻿"   # Arabic Presentation Forms-B
        for ch in text
    )


def _add_text_overlay(prs, slide, spec: dict) -> None:
    """An editable textbox positioned over a rendered slide background.

    Sized and placed from page fractions, with autofit off: the box is already
    the size the design gave it, and letting PowerPoint reflow it would move
    the text off its intended spot.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Emu, Pt

    # A little slack, because PDF bboxes hug the glyphs. It has to scale with
    # the TEXT, not with the slide: a flat 2% of slide height is 11.9 pt on
    # every box, which is a rounding error under a 38 pt headline but more
    # than a whole line-height under 10 pt body copy — so small text sat
    # visibly above where the design put it. pad_y is dropped entirely and the
    # box is grown downwards instead (see height below), which keeps the top
    # edge exactly where the PDF reported it.
    # Both are 0: a flat fraction of the SLIDE shifts every box by the same
    # amount whatever the type size, so it is a rounding error under a 58 pt
    # headline and a gross displacement under 10 pt body copy or an 18 pt
    # label (measured 11.9 pt up and 23.0 pt left). The box keeps the exact
    # top-left the PDF reported and is grown right and down instead.
    pad_x = 0
    pad_y = 0
    left = int(spec["x"] * prs.slide_width) - pad_x
    top = int(spec["y"] * prs.slide_height) - pad_y
    # Extra room RIGHT and BELOW only, so a wider substituted font or
    # PowerPoint's own line spacing cannot clip the text, while the top-left
    # corner stays exactly where the design put it.
    width = int(spec["w"] * prs.slide_width * 1.12) + int(prs.slide_width * 0.004)
    height = int(spec["h"] * prs.slide_height * 1.25)
    left = max(0, min(left, prs.slide_width - 1))
    top = max(0, min(top, prs.slide_height - 1))
    # These are EMU, so compare against EMU minimums — Pt() here fed points
    # into an EMU slot and raised "value must be in range 100 to 400000".
    _MIN_W = 152400   # 12 pt in EMU
    _MIN_H = 101600   # 8 pt
    width = max(_MIN_W, min(width, prs.slide_width - left))
    height = max(_MIN_H, min(height, prs.slide_height - top))

    if spec.get("align") == "center":
        centre = int((spec["x"] + spec["w"] / 2) * prs.slide_width)
        half = max(width // 2, int(prs.slide_width * 0.02))
        left = max(0, centre - half)
        width = min(2 * half, prs.slide_width - left)

    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    frame = box.text_frame
    frame.word_wrap = True
    try:
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    # Point sizes are relative to the source page height, so a 54 pt headline
    # on a 1080 pt page stays proportionally correct on the slide. BOTH sides
    # must be in points: slide_height is EMU, and dividing it by a point value
    # gave a scale of ~4762 and a 257,000 pt font.
    page_h_pt = float(spec.get("page_h") or 0)
    slide_h_pt = prs.slide_height / 12700.0        # EMU -> pt
    scale = (slide_h_pt / page_h_pt) if page_h_pt > 0 else 1.0

    align = spec.get("align", "left")
    for i, line in enumerate(spec["lines"]):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        if line.get("space_before") and i > 0:
            # space_before rather than an empty paragraph: an empty one takes
            # the run's full point size and overshoots a large headline.
            try:
                para.space_before = Pt(max(2.0, line["runs"][0]["size"] * scale * 0.6))
            except Exception:
                pass
        if align == "center":
            para.alignment = PP_ALIGN.CENTER
        # Arabic/Hebrew: PDF stores glyphs in visual order, so the extracted
        # string is already correct, but the paragraph must be marked RTL or
        # PowerPoint re-orders punctuation and digits to the wrong side.
        if any(_is_rtl(r["text"]) for r in line["runs"]):
            try:
                para._pPr.set("rtl", "1")
                if align != "center":
                    para.alignment = PP_ALIGN.RIGHT
            except Exception:
                pass
        for run_spec in line["runs"]:
            run = para.add_run()
            run.text = run_spec["text"]
            font = run.font
            pt_size = max(6.0, run_spec["size"] * scale)
            font.size = Pt(pt_size)
            # A single safe family rather than the PDF's embedded font: the
            # original is almost never installed, so PowerPoint substitutes
            # per-machine and the layout shifts. Calibri ships with Office on
            # every platform and covers Arabic and other non-Latin scripts.
            font.name = _OVERLAY_FONT
            font.bold = run_spec["bold"]
            font.italic = run_spec["italic"]
            # Letter spacing. DrawingML carries it as rPr/@spc in 1/100 pt;
            # python-pptx has no wrapper, so set the attribute directly.
            track = float(run_spec.get("track") or 0.0)
            if track > 0:
                try:
                    font._rPr.set("spc", str(int(round(track * pt_size * 100))))
                except Exception:
                    pass
            colour = run_spec["color"]
            font.color.rgb = RGBColor((colour >> 16) & 0xFF,
                                      (colour >> 8) & 0xFF, colour & 0xFF)


def _add_caption(prs, slide, text: str) -> None:
    """Small centred caption strip along the bottom of a picture slide.

    A TITLE_ONLY layout has no body placeholder, so a caption needs its own
    textbox or it would be silently dropped.
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt

    left = Emu(int(prs.slide_width * 0.06))
    top = Emu(int(prs.slide_height * 0.90))
    width = Emu(int(prs.slide_width * 0.88))
    height = Emu(int(prs.slide_height * 0.07))
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.italic = True


def _build_pptx_from_blocks(blocks: list[DocumentBlock], output_path: str) -> ConversionSummary:
    """Assemble a Presentation from extracted DocumentBlocks.

    Pictures — embedded images and rasterised vector figures alike — get their
    own slide, fitted to the content area with the aspect ratio preserved.
    """
    from pptx import Presentation
    from pptx.util import Pt

    summary = ConversionSummary()
    prs = Presentation()
    groups = _group_blocks_into_slides(blocks)

    # Default template slides are 4:3. A 16:9 source would be cropped on both
    # sides, so adopt the source page's aspect ratio when it is a design deck.
    _match_slide_size(prs, blocks)

    if not groups:
        # An empty deck is invalid to some readers; leave one explanatory slide.
        slide = prs.slides.add_slide(prs.slide_layouts[_PPTX_LAYOUT_TITLE_ONLY])
        slide.shapes.title.text = "No text content found"
        prs.save(output_path)
        return summary

    for group in groups:
        if group.get("full_bleed"):
            slide = prs.slides.add_slide(prs.slide_layouts[_PPTX_LAYOUT_BLANK])
            for block in group["blocks"]:
                if _add_full_bleed_picture(prs, slide, block):
                    summary.images += 1
                else:
                    summary.skipped_elements.append(
                        "Page %d could not be rendered" % (block.page_num + 1))
                for spec in getattr(block, "text_overlays", None) or []:
                    _add_text_overlay(prs, slide, spec)
                    summary.text_blocks += 1
            continue

        has_table = any(b.block_type == "table" for b in group["blocks"])
        has_image = any(b.block_type in ("image", "scanned_page")
                        for b in group["blocks"])
        layout = prs.slide_layouts[
            _PPTX_LAYOUT_TITLE_ONLY if (has_table or has_image)
            else _PPTX_LAYOUT_TITLE_CONTENT]
        slide = prs.slides.add_slide(layout)

        title_text = group["title"] or ""
        if group["continued"] and title_text:
            title_text += " (cont.)"
        if slide.shapes.title is not None:
            slide.shapes.title.text = title_text
            if title_text:
                summary.headings += 1

        if has_table:
            for block in group["blocks"]:
                if block.table_data:
                    _add_table_slide(prs, slide, block.table_data)
                    summary.tables += 1
            continue

        if has_image:
            for block in group["blocks"]:
                if block.block_type in ("image", "scanned_page"):
                    if _add_picture_slide(prs, slide, block):
                        if block.block_type == "scanned_page":
                            summary.scanned_pages += 1
                        else:
                            summary.images += 1
                    else:
                        summary.skipped_elements.append(
                            "Image on page %d" % (block.page_num + 1))
            caption = group.get("caption")
            if caption:
                _add_caption(prs, slide, caption)
                summary.text_blocks += 1
            continue

        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is None or not group["blocks"]:
            continue

        frame = body.text_frame
        frame.word_wrap = True
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            # Shrink rather than overflow if the estimate was optimistic.
            frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

        first = True
        for block in group["blocks"]:
            para = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            text = (block.text or "").strip()
            if block.block_type == "list_item":
                para.level = min(max((block.level or 1) - 1, 0), 4)
                summary.list_items += 1
            else:
                para.level = 0
                summary.text_blocks += 1
            # Runs carry the source's bold/italic where the extractor found it.
            if block.spans:
                for span in block.spans:
                    if not span.text:
                        continue
                    run = para.add_run()
                    run.text = span.text
                    run.font.bold = bool(span.bold)
                    run.font.italic = bool(span.italic)
                    run.font.size = Pt(18)
            else:
                run = para.add_run()
                run.text = text
                run.font.size = Pt(18)

    # Images are placed above; only note pages whose text could not be read.
    for block in blocks:
        if block.block_type == "scanned_page":
            summary.warnings.append(
                "Page %d is a scanned image - its text is not searchable."
                % (block.page_num + 1))

    prs.save(output_path)
    return summary


def _pdf_is_design_deck(doc) -> bool:
    """True when the PDF as a whole is a composed presentation.

    Sampled across the document rather than judged per page: a deck usually has
    a few pages without a background image (a table, a divider), and rendering
    35 pages whole while re-flowing one into 22 slides is worse than treating
    the file consistently.
    """
    pages = len(doc)
    if pages == 0:
        return False
    step = max(1, pages // 12)
    checked = 0
    designed = 0
    for i in range(0, pages, step):
        checked += 1
        if _page_is_designed_slide(doc[i]):
            designed += 1
    return checked > 0 and designed / checked >= 0.6


def _pdf_to_pptx(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """PDF to PowerPoint.

    A composed design deck is rendered page-for-slide, preserving its layout.
    A flowed document is decomposed into text blocks and packed into slides.
    """
    try:
        with fitz.open(input_path) as doc:
            if doc.is_encrypted:
                return False, "PDF is encrypted/password-protected", None
            total_pages = len(doc)
            is_deck = _pdf_is_design_deck(doc)
            _log("PPTX: %s" % ("design deck — rendering page-for-slide"
                               if is_deck else "document — rebuilding slides"))
            body_font_size = _detect_body_font_size(doc)
            all_blocks = []
            extracted_images = set()
            warnings: list[str] = []
            for page_num, page in enumerate(doc):
                if cancel_event and cancel_event.is_set():
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    return False, "Conversion cancelled by user", None
                if progress_callback:
                    progress_callback(page_num + 1, total_pages)
                # In a design deck EVERY page is rendered whole, including the
                # occasional page without a background image — mixing the two
                # strategies inside one file gives wildly inconsistent output.
                all_blocks.extend(_extract_page_blocks(
                    page, page_num, body_font_size, doc, extracted_images,
                    warnings=warnings,
                    whole_slide=is_deck or _page_is_designed_slide(page)))
            summary = _build_pptx_from_blocks(all_blocks, output_path)
            summary.total_pages = total_pages
            if warnings:
                summary.warnings.extend(warnings)
            return True, output_path, summary
    except Exception as e:
        _log("PDF to PPTX failed: %s" % e)
        return False, str(e), None


def _docx_to_pptx(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """DOCX to PowerPoint, text only."""
    try:
        blocks = _docx_to_blocks(input_path, cancel_event)
        if blocks is None:
            return False, "Conversion cancelled by user", None
        if progress_callback:
            progress_callback(1, 1)
        summary = _build_pptx_from_blocks(blocks, output_path)
        summary.total_pages = 1
        return True, output_path, summary
    except Exception as e:
        _log("DOCX to PPTX failed: %s" % e)
        return False, str(e), None


def _pptx_to_blocks(input_path: str, progress_callback=None, cancel_event=None) -> list[DocumentBlock]:
    """Text of every slide as DocumentBlocks. Pictures and charts are ignored.

    A slide's title placeholder becomes an H2 and its body text frame becomes
    paragraphs or list items, using PowerPoint's own outline level for nesting.
    Tables are read; images, charts and SmartArt are skipped by design.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(input_path)
    blocks: list[DocumentBlock] = []
    slides = list(prs.slides)

    for index, slide in enumerate(slides):
        if cancel_event and cancel_event.is_set():
            return []
        if progress_callback:
            progress_callback(index + 1, len(slides))

        # python-pptx builds a NEW proxy object on every `.title` access, so an
        # identity check against the iterated shapes never matches and the title
        # gets emitted twice (once as a heading, once as a bullet). Compare the
        # underlying shape_id instead.
        title_id = None
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title_id = title_shape.shape_id
                title = title_shape.text_frame.text.strip()
                if title:
                    blocks.append(DocumentBlock("heading", text=title, level=2,
                                                page_num=index))
        except Exception:
            title_id = None

        for shape in slide.shapes:
            if title_id is not None and shape.shape_id == title_id:
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text for r in para.runs).strip()
                    if not text:
                        continue
                    spans = [SpanInfo(text=r.text, bold=bool(r.font.bold),
                                      italic=bool(r.font.italic), font_size=0.0,
                                      font_name="", color=(0, 0, 0))
                             for r in para.runs if r.text]
                    # PowerPoint bullets live in the theme, not the text, so
                    # outline level is the only reliable list signal. Level 0
                    # in a body placeholder is still a bullet in practice.
                    level = (para.level or 0) + 1
                    blocks.append(DocumentBlock(
                        "list_item", text=text, spans=spans, level=level,
                        list_style="bullet", page_num=index))

            elif getattr(shape, "has_table", False) and shape.has_table:
                rows = [[c.text.strip() for c in r.cells]
                        for r in shape.table.rows]
                if rows:
                    blocks.append(DocumentBlock("table", table_data=rows,
                                                page_num=index))

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Recorded so the summary can say what was left out.
                blocks.append(DocumentBlock("image", page_num=index))

    return blocks


def _pptx_to_md(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """PPTX to Markdown, text only."""
    try:
        blocks = _pptx_to_blocks(input_path, progress_callback, cancel_event)
        if cancel_event and cancel_event.is_set():
            return False, "Conversion cancelled by user", None
        summary = _build_md_from_blocks(blocks, output_path)
        summary.total_pages = len({b.page_num for b in blocks}) or 1
        return True, output_path, summary
    except Exception as e:
        _log("PPTX to Markdown failed: %s" % e)
        return False, str(e), None


def _pptx_to_docx(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """PPTX to DOCX, text only."""
    try:
        blocks = _pptx_to_blocks(input_path, progress_callback, cancel_event)
        if cancel_event and cancel_event.is_set():
            return False, "Conversion cancelled by user", None
        summary = _build_docx_from_blocks(blocks, output_path)
        summary.total_pages = len({b.page_num for b in blocks}) or 1
        return True, output_path, summary
    except Exception as e:
        _log("PPTX to DOCX failed: %s" % e)
        return False, str(e), None


def _pdf_to_docx(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """Orchestrate high-fidelity PDF to DOCX conversion."""
    try:
        with fitz.open(input_path) as doc:
            if doc.is_encrypted:
                return False, "PDF is encrypted/password-protected", None
                
            total_pages = len(doc)
            body_font_size = _detect_body_font_size(doc)
            all_blocks = []
            extracted_images = set() # For deduplication (T016)
            extraction_warnings = []
            
            for page_num, page in enumerate(doc):
                # Check for cancellation (FR-014)
                if cancel_event and cancel_event.is_set():
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    return False, "Conversion cancelled by user", None
                    
                if progress_callback:
                    progress_callback(page_num + 1, total_pages)
                
                page_blocks = _extract_page_blocks(page, page_num, body_font_size, doc, extracted_images, warnings=extraction_warnings)
                all_blocks.extend(page_blocks)
                
            summary = _build_docx_from_blocks(all_blocks, output_path)
            summary.total_pages = total_pages
            if extraction_warnings:
                summary.warnings.extend(extraction_warnings)
            
            return True, output_path, summary
            
    except Exception as e:
        _log(f"PDF to DOCX failed: {e}")
        return False, str(e), None


import threading
_libreoffice_lock = threading.Lock()

def detect_converter_backend() -> str:
    """Detect available DOCX to PDF converter backend."""
    import sys
    import shutil

    docx2pdf_available = False
    try:
        import docx2pdf  # noqa: F401
        docx2pdf_available = True
    except ImportError:
        pass

    if docx2pdf_available and sys.platform == "win32":
        # Check Word is installed via registry (avoids launching Word as a side-effect)
        try:
            import winreg
            winreg.OpenKey(
                winreg.HKEY_LOCAL_MACHINE,
                r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\WINWORD.EXE",
            )
            return "word"
        except Exception:
            pass
        # Fallback: try registry under Wow6432Node (32-bit Word on 64-bit OS)
        try:
            import winreg
            winreg.OpenKey(
                winreg.HKEY_LOCAL_MACHINE,
                r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\WINWORD.EXE",
            )
            return "word"
        except Exception:
            pass

    if docx2pdf_available and sys.platform == "darwin":  # type: ignore[misc]
        # macOS: docx2pdf drives Word via AppleScript
        if os.path.exists("/Applications/Microsoft Word.app"):
            return "word"

    if shutil.which("libreoffice") or shutil.which("soffice"):
        return "libreoffice"

    if sys.platform == "darwin" and os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        return "libreoffice"

    return "none"

def _docx_to_pdf(input_path: str, output_path: str, progress_callback=None, cancel_event=None) -> tuple[bool, str, ConversionSummary | None]:
    """Convert DOCX to PDF using docx2pdf or LibreOffice fallback."""
    import sys
    import shutil
    import subprocess
    
    backend = detect_converter_backend()
    if backend == "none":
        return False, "No DOCX-to-PDF converter found (install Microsoft Word or LibreOffice)", None
        
    try:
        from docx import Document
        doc = Document(input_path)
        para_count = len(doc.paragraphs)
    except Exception:
        para_count = 0
    summary = ConversionSummary(total_pages=1, text_blocks=para_count)
        
    if backend == "word":
        try:
            from docx2pdf import convert
            _log("Using docx2pdf backend...")
            convert(input_path, output_path)
            if os.path.exists(output_path):
                return True, output_path, summary
        except Exception as e:
            _log(f"docx2pdf failed, falling back to LibreOffice if available: {e}")
            pass

    soffice_cmd = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice_cmd and sys.platform == "darwin":
        if os.path.exists("/Applications/LibreOffice.app/Contents/MacOS/soffice"):
            soffice_cmd = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            
    if soffice_cmd:
        _log("Using LibreOffice backend...")
        outdir = os.path.dirname(os.path.abspath(output_path))
        abs_input = os.path.abspath(input_path)
        cmd = [
            soffice_cmd,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", outdir,
            abs_input
        ]
        with _libreoffice_lock:
            try:
                kwargs = {"creationflags": 0x08000000} if sys.platform == "win32" else {}
                
                # Check cancellation before starting subprocess
                if cancel_event and cancel_event.is_set():
                    return False, "Conversion cancelled by user", None
                    
                # We can't interrupt the subprocess cleanly cross-platform here without complex logic,
                # but we poll cancel_event in convert_document wrapper or rely on subprocess timeouts.
                result = subprocess.run(cmd, capture_output=True, timeout=300, **kwargs)
                if result.returncode == 0:
                    expected_out = os.path.join(outdir, os.path.splitext(os.path.basename(abs_input))[0] + ".pdf")
                    abs_output = os.path.abspath(output_path)
                    if expected_out != abs_output and os.path.exists(expected_out):
                        if os.path.exists(abs_output):
                            os.remove(abs_output)
                        shutil.move(expected_out, abs_output)
                    if os.path.exists(abs_output):
                        return True, output_path, summary
                return False, f"LibreOffice conversion failed: {result.stderr.decode()}", None
            except Exception as e:
                return False, f"LibreOffice failed: {e}", None
                
    return False, "Conversion failed (no backend available)", None


def convert_document(
    input_path: str, 
    output_format: str, 
    progress_callback=None, 
    cancel_event=None
) -> tuple[bool, str, ConversionSummary | None]:
    """Convert a document to output_format.

    Supports PDF and DOCX as input, and the same set as output.
    Images (JPG/PNG/BMP/GIF/WEBP) can be converted to PDF.

    Returns: (success, output_path_or_error, summary)
    """
    input_ext = os.path.splitext(input_path)[1].lower()
    base_name  = os.path.splitext(input_path)[0]
    output_path = f"{base_name}_converted.{output_format}"

    if output_format not in _VALID_OUTPUTS:
        return False, f"Unsupported output format: {output_format}", None
    if input_ext not in _VALID_DOCS | _IMAGE_EXTS:
        return False, f"Unsupported input format: {input_ext}", None

    _log(f"Starting: {os.path.basename(input_path)}  →  {output_format.upper()}")

    # ------------------------------------------------------------------
    # PDF → other
    # ------------------------------------------------------------------
    if input_ext == ".pdf":
        if output_format == "docx":
            return _pdf_to_docx(input_path, output_path, progress_callback, cancel_event)
        if output_format == "md":
            return _pdf_to_md(input_path, output_path, progress_callback, cancel_event)
        if output_format == "pptx":
            return _pdf_to_pptx(input_path, output_path, progress_callback, cancel_event)
        return False, f"PDF → {output_format} is not supported", None

    # ------------------------------------------------------------------
    # DOCX → other
    # ------------------------------------------------------------------
    elif input_ext == ".docx":
        if output_format == "pdf":
            return _docx_to_pdf(input_path, output_path, progress_callback, cancel_event)
        if output_format == "md":
            return _docx_to_md(input_path, output_path, progress_callback, cancel_event)
        if output_format == "pptx":
            return _docx_to_pptx(input_path, output_path, progress_callback, cancel_event)
        return False, f"DOCX → {output_format} is not supported", None

    # ------------------------------------------------------------------
    # Image → PDF
    # ------------------------------------------------------------------
    # ------------------------------------------------------------------
    # PPTX → other (text only; pictures/charts are ignored by design)
    # ------------------------------------------------------------------
    elif input_ext == ".pptx":
        if output_format == "md":
            return _pptx_to_md(input_path, output_path, progress_callback, cancel_event)
        if output_format == "docx":
            return _pptx_to_docx(input_path, output_path, progress_callback, cancel_event)
        return False, f"PPTX → {output_format} is not supported", None

    # ------------------------------------------------------------------
    # Image → PDF
    # ------------------------------------------------------------------
    elif input_ext in _IMAGE_EXTS:
        if output_format != "pdf":
            return False, f"Image → {output_format} is not supported (only image → PDF)", None
        _log("Converting image to PDF…")
        summary = ConversionSummary(total_pages=1, images=1)
        try:
            with Image.open(input_path) as image:
                _log(f"  Image size: {image.size}, mode: {image.mode}")
                img_bytes = io.BytesIO()
                image.convert("RGB").save(img_bytes, format="PNG")
        except Exception as e:
            return False, f"Error reading image: {e}", None
            
        pdf_doc = fitz.open()
        try:
            page = pdf_doc.new_page()
            page.insert_image(page.rect, stream=img_bytes.getvalue())
            _log(f"Saving PDF → {output_path}")
            try:
                pdf_doc.save(output_path)
                return True, output_path, summary
            except Exception as e:
                return False, f"ERROR saving PDF: {e}", None
        finally:
            pdf_doc.close()
