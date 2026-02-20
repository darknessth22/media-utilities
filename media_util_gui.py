import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os
import sys
import subprocess
from yt_dlp import YoutubeDL
from PIL import Image, ExifTags
import threading
import fitz  # PyMuPDF
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import io
# glitch
def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])
def check_dependencies():
    """Check all runtime dependencies. Returns 'ffmpeg_missing' if FFmpeg is absent, else None."""
    try:
        import yt_dlp
    except ImportError:
        print("Installing yt-dlp...")
        install("yt-dlp")

    try:
        subprocess.run([ffmpeg_path, "-version"], capture_output=True, check=True)
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("\nFFmpeg not found! Media conversion features will be unavailable.")
        return "ffmpeg_missing"

    try:
        import pillow_heif
    except ImportError:
        print("\nInstalling pillow-heif for HEIC support...")
        install("pillow-heif")
        import pillow_heif
        pillow_heif.register_heif_opener()

    # spotdl (Spotify support) — optional, not available on Python 3.14+
    try:
        result = subprocess.run(["spotdl", "--version"], capture_output=True, text=True, check=True)
        print(f"spotdl available: {result.stdout.strip()}")
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("spotdl not found — Spotify downloads will be unavailable.")
        print("Install on Python <=3.13 with: pip install 'spotdl>=3.9.6'")

    return None
import sys
import os

def _find_binary(name):
    """Locate a binary using a four-step fallback chain."""
    import shutil
    exe_name = f"{name}.exe" if sys.platform == "win32" else name
    # 1. PyInstaller bundle (onefile or onedir)
    if hasattr(sys, '_MEIPASS'):
        bundled = os.path.join(sys._MEIPASS, exe_name)
        if os.path.isfile(bundled):
            return bundled
    # 2. Same directory as this script / frozen executable
    script_dir = os.path.dirname(os.path.abspath(__file__))
    local = os.path.join(script_dir, exe_name)
    if os.path.isfile(local):
        return local
    # 3. Current working directory (legacy behaviour)
    cwd = os.path.join(os.getcwd(), exe_name)
    if os.path.isfile(cwd):
        return cwd
    # 4. System PATH
    which = shutil.which(name)
    if which:
        return which
    # Return bare name — callers receive a clear FileNotFoundError at call time
    return exe_name

def get_ffmpeg_path():
    return _find_binary("ffmpeg")

def get_ffprobe_path():
    return _find_binary("ffprobe")

ffmpeg_path = get_ffmpeg_path()
ffprobe_path = get_ffprobe_path()
print(f"FFmpeg path: {ffmpeg_path}")
print(f"FFprobe path: {ffprobe_path}")
def get_platform(url):
    domains = {
        'youtube': ['youtube.com', 'youtu.be'],
        'facebook': ['facebook.com', 'fb.watch'],
        'instagram': ['instagram.com', 'instagr.am'],
        'tiktok': ['tiktok.com'],
        'twitter': ['twitter.com', 'x.com'],
        'spotify': ['spotify.com', 'open.spotify.com']
    }
    for platform, urls in domains.items():
        if any(domain in url for domain in urls):
            return platform
    return 'generic'

def parse_time(time_str):
    parts = list(map(int, time_str.split(':')))
    if len(parts) == 1:
        return parts[0]
    elif len(parts) == 2:
        return parts[0] * 60 + parts[1]
    elif len(parts) == 3:
        return parts[0] * 3600 + parts[1] * 60 + parts[2]
    else:
        raise ValueError("Invalid time format. Use H:MM:SS, M:SS, or S")

def download_spotify(url, audio_format="mp3", output_dir=None):
    """Download Spotify tracks using spotdl (gets metadata from Spotify, audio from YouTube)"""
    try:
        # Build spotdl command
        cmd = ["spotdl", "download", url]

        # Set output directory if specified
        if output_dir:
            cmd.extend(["--output", output_dir])

        # Set audio format
        if audio_format in ["mp3", "flac", "ogg", "opus", "m4a"]:
            cmd.extend(["--format", audio_format])
        else:
            cmd.extend(["--format", "mp3"])  # Default to mp3

        # Set audio quality
        cmd.extend(["--bitrate", "320k"])

        # Add other useful options
        cmd.extend([
            "--threads", "4",  # Use 4 threads for faster download
            "--sponsor-block",  # Skip sponsor segments
        ])

        print(f"Running spotdl command: {' '.join(cmd)}")

        # Run spotdl
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)

        print("Spotify download completed successfully!")
        print("Note: Audio sourced from YouTube with Spotify metadata")

        return True

    except subprocess.CalledProcessError as e:
        print(f"spotdl error: {e}")
        if e.stderr:
            print(f"Error details: {e.stderr}")
        return False
    except Exception as e:
        print(f"Error downloading from Spotify: {str(e)}")
        return False

def download_media(url, platform, media_type='video', quality=None, start_time=None, end_time=None, audio_format="mp3", output_dir=None, video_codec="libx264", force_codec=False):
    # Handle Spotify URLs with spotdl
    if platform == 'spotify':
        print("Detected Spotify URL - using spotdl for download")
        print("Note: This will get metadata from Spotify and audio from YouTube")
        return download_spotify(url, audio_format, output_dir)

    # Set output template based on output directory
    output_template = os.path.join(output_dir, '%(title)s.%(ext)s') if output_dir else '%(title)s.%(ext)s'
    ydl_opts = {
        'outtmpl': output_template,
        'cookiefile': 'cookies.txt' if os.path.exists('cookies.txt') else None,
        'postprocessor_args': ['-loglevel', 'error'],
        'force_keyframes_at_cuts': True
    }

    if start_time and end_time:
        start = parse_time(start_time)
        end = parse_time(end_time)

        # Just set the output format to mp4 and add trimming options
        ydl_opts['merge_output_format'] = 'mp4'
        ydl_opts['postprocessor_args'] = [
            '-ss', str(start),         # Start time
            '-to', str(end)            # End time
        ]
        # Update output template for trimmed videos while preserving the output directory
        if output_dir:
            ydl_opts['outtmpl'] = os.path.join(output_dir, f'%(title)s_Trimmed_{start}s_{end}s.%(ext)s')
        else:
            ydl_opts['outtmpl'] = f'%(title)s_Trimmed_{start}s_{end}s.%(ext)s'

    if media_type == 'audio':
        ydl_opts['format'] = 'bestaudio/best'
        ydl_opts['postprocessors'] = [{
            'key': 'FFmpegExtractAudio',
            'preferredcodec': audio_format,
            'preferredquality': '320'
        }]
    else:
        ydl_opts['format'] = quality or 'bestvideo+bestaudio/best'

        # Just set the output format to mp4, we'll handle codec conversion separately
        ydl_opts['merge_output_format'] = 'mp4'

    try:
        # First, download the video with yt-dlp
        with YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            downloaded_file = ydl.prepare_filename(info)

        # If we're downloading video and not audio only, check and convert to h264 if needed
        if media_type != 'audio' and os.path.exists(downloaded_file):
            # Get detailed information about the video
            info_cmd = [
                ffprobe_path, '-v', 'error', '-select_streams', 'v:0',
                '-show_entries', 'stream=codec_name,profile,level,bit_rate', '-of', 'json',
                downloaded_file
            ]

            try:
                import json
                codec_info = json.loads(subprocess.check_output(info_cmd, text=True))
                stream_info = codec_info.get('streams', [{}])[0]
                current_codec = stream_info.get('codec_name', '').lower()
                profile = stream_info.get('profile', '').lower()
                bit_rate = stream_info.get('bit_rate', 'unknown')

                print(f"Video details: codec={current_codec}, profile={profile}, bitrate={bit_rate}")

                # Check if it's h264/AVC
                is_h264 = current_codec in ['h264', 'avc', 'avc1']

                # If not h264, or if force_codec is True, convert it
                if not is_h264 or force_codec:
                    print(f"Converting from {current_codec} to h264...")
                    # Create a temporary filename for the re-encoded file
                    base, ext = os.path.splitext(downloaded_file)
                    temp_file = f"{base}_h264{ext}"

                    # Re-encode with FFmpeg to ensure h264 codec
                    cmd = [
                        'ffmpeg', '-y',
                        '-i', downloaded_file,
                        '-c:v', video_codec,  # Use h264 codec
                        '-preset', 'fast',    # Fast encoding preset
                        '-crf', '23',         # Quality level (lower is better)
                        '-profile:v', 'main', # Use main profile for better compatibility
                        '-level', '4.0',      # Set level for compatibility
                        '-pix_fmt', 'yuv420p', # Use standard pixel format
                        '-c:a', 'copy',       # Copy audio to preserve quality
                        temp_file
                    ]

                    subprocess.run(cmd, check=True, capture_output=True)

                    # Replace the original file with the re-encoded one
                    if os.path.exists(temp_file):
                        os.remove(downloaded_file)
                        os.rename(temp_file, downloaded_file)
                        print("Successfully converted to h264")
                else:
                    # Even if it's h264, we might want to force a specific profile or bitrate
                    # For now, we'll just print a message, but you could add conditions here
                    # to re-encode based on profile or bitrate if needed
                    print("Video is already using h264/AVC codec, no conversion needed")

                    # Uncomment the following lines if you want to force re-encoding regardless
                    # of the current codec (e.g., to ensure a specific profile or bitrate)
                    '''
                    print("Force re-encoding to ensure specific h264 settings...")
                    base, ext = os.path.splitext(downloaded_file)
                    temp_file = f"{base}_h264{ext}"

                    cmd = [
                        'ffmpeg', '-y',
                        '-i', downloaded_file,
                        '-c:v', video_codec,
                        '-preset', 'fast',
                        '-crf', '23',
                        '-c:a', 'copy',
                        temp_file
                    ]

                    subprocess.run(cmd, check=True, capture_output=True)

                    if os.path.exists(temp_file):
                        os.remove(downloaded_file)
                        os.rename(temp_file, downloaded_file)
                        print("Successfully re-encoded to h264 with specific settings")
                    '''
            except Exception as e:
                print(f"Error checking codec: {e}")
                # If we can't check the codec, try to convert anyway
                print("Attempting conversion without codec information...")
                base, ext = os.path.splitext(downloaded_file)
                temp_file = f"{base}_h264{ext}"

                cmd = [
                    'ffmpeg', '-y',
                    '-i', downloaded_file,
                    '-c:v', video_codec,  # Use h264 codec
                    '-preset', 'fast',    # Fast encoding preset
                    '-crf', '23',         # Quality level (lower is better)
                    '-profile:v', 'main', # Use main profile for better compatibility
                    '-level', '4.0',      # Set level for compatibility
                    '-pix_fmt', 'yuv420p', # Use standard pixel format
                    '-c:a', 'copy',       # Copy audio to preserve quality
                    temp_file
                ]

                subprocess.run(cmd, check=True, capture_output=True)

                if os.path.exists(temp_file):
                    os.remove(downloaded_file)
                    os.rename(temp_file, downloaded_file)
                    print("Conversion completed")

        return True
    except Exception as e:
        print(f"Error downloading: {str(e)}")
        return False

def convert_images(file_paths, target_format):
    try:
        import pillow_heif
        pillow_heif.register_heif_opener()
    except ImportError:
        install("pillow-heif")
        import pillow_heif
        pillow_heif.register_heif_opener()

    format_mapping = {
        "jpg": "JPEG",
        "jpeg": "JPEG",
        "png": "PNG",
        "bmp": "BMP",
        "gif": "GIF",
        "webp": "WEBP",
        "heic": "HEIF",
        "heif": "HEIF"
    }

    target_format = format_mapping.get(target_format.lower(), target_format.upper())
    if target_format not in format_mapping.values():
        print(f"Unsupported target format: {target_format}")
        return False

    success_count = 0
    for file_path in file_paths:
        try:
            image = Image.open(file_path)
            exif_data = image.info.get("exif")
            output_path = file_path.rsplit(".", 1)[0] + f".{target_format.lower()}"

            if target_format == "HEIF":
                heif_options = {
                    "quality": 90,
                    "compression": "hevc",
                    "bit_depth": 8,
                    "chroma_subsampling": "420",
                    "save_all": False,
                }
                image.save(output_path, format="HEIF", exif=exif_data, **heif_options)
            else:
                image.save(output_path, format=target_format, quality=95)

            success_count += 1
        except Exception as e:
            print(f"Conversion failed for {file_path}: {e}")

    return success_count > 0

def convert_media(file_path):
    if not os.path.exists(file_path):
        return False

    supported_formats = {
        'audio': {'mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a'},
        'video': {'mp4', 'mkv', 'avi', 'mov', 'webm', 'flv'},
        'image': {'jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif'}
    }

    ext = os.path.splitext(file_path)[1][1:].lower()
    media_type = None
    for category, exts in supported_formats.items():
        if ext in exts:
            media_type = category
            break

    if not media_type:
        return False

    if media_type == "image":
        return convert_images([file_path], ext)

    base = os.path.splitext(file_path)[0]
    output_path = f"{base}_converted.{ext}"

    try:
        cmd = [ffmpeg_path, '-y', '-i', file_path, output_path]
        subprocess.run(cmd, check=True, capture_output=True)
        return True
    except subprocess.CalledProcessError:
        return False

def convert_document(input_path, output_format):
    """Convert documents between different formats with enhanced layout and formatting preservation"""
    input_ext = os.path.splitext(input_path)[1].lower()
    base_name = os.path.splitext(input_path)[0]
    output_path = f"{base_name}_converted.{output_format}"

    # PDF to other formats
    if input_ext == '.pdf':
        doc = fitz.open(input_path)

        if output_format == 'docx':
            word_doc = Document()

            for page_num, page in enumerate(doc):
                # Get structured text with formatting information
                text_dict = page.get_text("dict")

                # Get image information with positions
                image_list = page.get_images()
                image_rects = []

                # Extract image data (simplified approach)
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)

                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")

                            # Store image info (without complex positioning for now)
                            image_rects.append({
                                'data': img_data,
                                'index': img_index,
                                'page': page_num,
                                'width': pix.width,
                                'height': pix.height
                            })

                        pix = None  # Free memory
                    except Exception as e:
                        print(f"Warning: Could not process image {img_index} from page {page_num}: {e}")
                        continue

                # Process text blocks with formatting
                for block in text_dict.get("blocks", []):
                    if "lines" in block:  # Text block
                        block_text = ""
                        block_font_size = 12  # Default
                        block_alignment = None

                        for line in block["lines"]:
                            line_text = ""
                            for span in line.get("spans", []):
                                text = span.get("text", "")
                                font_size = span.get("size", 12)
                                font_flags = span.get("flags", 0)

                                # Update block font size (use largest in block)
                                if font_size > block_font_size:
                                    block_font_size = font_size

                                line_text += text

                            if line_text.strip():
                                block_text += line_text + "\n"

                        if block_text.strip():
                            # Determine alignment based on block position
                            block_rect = fitz.Rect(block["bbox"])
                            page_width = page.rect.width

                            # Simple alignment detection
                            if block_rect.x0 > page_width * 0.4 and block_rect.x1 < page_width * 0.6:
                                block_alignment = 'center'
                            elif block_rect.x0 > page_width * 0.7:
                                block_alignment = 'right'
                            else:
                                block_alignment = 'left'

                            # Add paragraph with formatting
                            paragraph = word_doc.add_paragraph(block_text.strip())

                            # Apply alignment
                            from docx.enum.text import WD_ALIGN_PARAGRAPH
                            if block_alignment == 'center':
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            elif block_alignment == 'right':
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            else:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

                            # Apply font size (approximate conversion)
                            for run in paragraph.runs:
                                from docx.shared import Pt
                                run.font.size = Pt(max(8, min(block_font_size, 24)))

                        # Add images after each text block (simplified approach)
                        # Process one image per text block if available
                        if image_rects and not any(img.get('processed', False) for img in image_rects):
                            # Get the first unprocessed image
                            for img_info in image_rects:
                                if not img_info.get('processed', False):
                                    try:
                                        # Create temporary image file
                                        temp_img_path = f"temp_img_{img_info['page']}_{img_info['index']}.png"
                                        with open(temp_img_path, "wb") as img_file:
                                            img_file.write(img_info['data'])

                                        # Calculate appropriate size (maintain aspect ratio)
                                        from docx.shared import Inches
                                        original_width = img_info['width']
                                        original_height = img_info['height']

                                        # Scale to fit page width (max 6 inches)
                                        max_width = 6.0
                                        if original_width > 0:
                                            scale_factor = min(max_width * 72 / original_width, 1.0)
                                            img_width = Inches(original_width * scale_factor / 72)
                                            img_height = Inches(original_height * scale_factor / 72)
                                        else:
                                            img_width = Inches(4)
                                            img_height = None

                                        # Add image to document
                                        img_paragraph = word_doc.add_paragraph()
                                        run = img_paragraph.add_run()

                                        if img_height:
                                            run.add_picture(temp_img_path, width=img_width, height=img_height)
                                        else:
                                            run.add_picture(temp_img_path, width=img_width)

                                        # Center align images by default
                                        from docx.enum.text import WD_ALIGN_PARAGRAPH
                                        img_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

                                        # Clean up temporary file
                                        os.remove(temp_img_path)

                                        # Mark image as processed
                                        img_info['processed'] = True
                                        break  # Process only one image per text block

                                    except Exception as e:
                                        print(f"Warning: Could not add image {img_info['index']}: {e}")
                                        break

                # Add any remaining unprocessed images at the end of the page
                for img_info in image_rects:
                    if not img_info.get('processed', False):
                        try:
                            temp_img_path = f"temp_img_{img_info['page']}_{img_info['index']}.png"
                            with open(temp_img_path, "wb") as img_file:
                                img_file.write(img_info['data'])

                            from docx.shared import Inches
                            paragraph = word_doc.add_paragraph()
                            run = paragraph.add_run()

                            # Calculate size maintaining aspect ratio
                            original_width = img_info['width']
                            original_height = img_info['height']
                            if original_width > 0:
                                scale_factor = min(6.0 * 72 / original_width, 1.0)
                                img_width = Inches(original_width * scale_factor / 72)
                                img_height = Inches(original_height * scale_factor / 72)
                                run.add_picture(temp_img_path, width=img_width, height=img_height)
                            else:
                                run.add_picture(temp_img_path, width=Inches(4))

                            os.remove(temp_img_path)

                        except Exception as e:
                            print(f"Warning: Could not add remaining image {img_info['index']}: {e}")

                # Add page break if not the last page
                if page_num < len(doc) - 1:
                    word_doc.add_page_break()

            word_doc.save(output_path)

        elif output_format == 'pptx':
            prs = Presentation()

            for page_num, page in enumerate(doc):
                # Create new slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

                # Add text if present
                text = page.get_text()
                if text.strip():
                    # Add text box for content
                    from pptx.util import Inches
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = prs.slide_width - Inches(1)
                    height = Inches(2)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = text

                # Add images
                image_list = page.get_images()
                img_top = Inches(3)  # Start images below text

                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)

                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")

                            # Create temporary image file
                            temp_img_path = f"temp_img_{page_num}_{img_index}.png"
                            with open(temp_img_path, "wb") as img_file:
                                img_file.write(img_data)

                            # Add image to slide
                            left = Inches(1)
                            width = Inches(6)
                            slide.shapes.add_picture(temp_img_path, left, img_top, width=width)

                            # Move next image down
                            img_top += Inches(2)

                            # Clean up temporary file
                            os.remove(temp_img_path)

                        pix = None  # Free memory
                    except Exception as e:
                        print(f"Warning: Could not add image {img_index} to slide {page_num}: {e}")
                        continue

            prs.save(output_path)

        elif output_format == 'xlsx':
            wb = Workbook()
            ws = wb.active
            for i, page in enumerate(doc):
                text = page.get_text()
                ws.cell(row=i+1, column=1, value=text)
            wb.save(output_path)

    # Word to other formats
    elif input_ext == '.docx':
        doc = Document(input_path)

        if output_format == 'pdf':
            # Enhanced Word to PDF conversion with proper image extraction and formatting
            pdf_doc = fitz.open()
            page = pdf_doc.new_page()
            y_position = 50
            page_width = page.rect.width
            margin = 50

            for paragraph in doc.paragraphs:
                # Handle text with formatting and alignment
                if paragraph.text.strip():
                    # Determine alignment
                    alignment = paragraph.alignment
                    x_position = margin

                    # Calculate text width for alignment
                    text_width = len(paragraph.text) * 6  # Approximate character width

                    if alignment == 1:  # Center
                        x_position = (page_width - text_width) / 2
                    elif alignment == 2:  # Right
                        x_position = page_width - margin - text_width

                    # Ensure x_position is within bounds
                    x_position = max(margin, min(x_position, page_width - margin - 100))

                    # Get font size from first run if available
                    font_size = 12
                    if paragraph.runs:
                        first_run = paragraph.runs[0]
                        if first_run.font.size:
                            font_size = min(24, max(8, first_run.font.size.pt))

                    # Insert text with proper positioning
                    text_rect = fitz.Rect(x_position, y_position, page_width - margin, y_position + font_size + 5)
                    page.insert_text(text_rect.tl, paragraph.text, fontsize=font_size)
                    y_position += font_size + 8

                # Extract and add actual images from paragraph runs
                for run in paragraph.runs:
                    if run.element.xml:
                        import xml.etree.ElementTree as ET
                        try:
                            # Parse the XML to find image relationships
                            root = ET.fromstring(run.element.xml)

                            # Look for drawing elements with image references
                            for elem in root.iter():
                                if 'blip' in str(elem.tag).lower() and 'embed' in elem.attrib:
                                    # Found an image reference
                                    rel_id = elem.attrib['embed']

                                    # Get the image from document relationships
                                    try:
                                        # Access document relationships to get actual image
                                        from docx.opc.constants import RELATIONSHIP_TYPE as RT

                                        # Find the image part
                                        for rel in doc.part.rels.values():
                                            if rel.rId == rel_id:
                                                image_part = rel.target_part
                                                img_bytes = image_part.blob

                                                # Calculate image size and position
                                                img_width = 200  # Default width
                                                img_height = 150  # Default height

                                                # Try to get original dimensions from the drawing
                                                for extent in root.iter():
                                                    if 'extent' in str(extent.tag).lower():
                                                        if 'cx' in extent.attrib and 'cy' in extent.attrib:
                                                            # Convert EMU to points (1 EMU = 1/914400 inch, 1 inch = 72 points)
                                                            img_width = min(400, int(extent.attrib['cx']) / 914400 * 72)
                                                            img_height = min(300, int(extent.attrib['cy']) / 914400 * 72)
                                                            break

                                                # Center image horizontally
                                                img_x = (page_width - img_width) / 2
                                                img_rect = fitz.Rect(img_x, y_position, img_x + img_width, y_position + img_height)

                                                # Insert image into PDF
                                                page.insert_image(img_rect, stream=img_bytes)
                                                y_position += img_height + 10
                                                break
                                    except Exception as img_e:
                                        # If image extraction fails, add a placeholder
                                        page.insert_text((margin, y_position), "[Image - could not extract]", fontsize=10)
                                        y_position += 20
                                        print(f"Warning: Could not extract image: {img_e}")

                        except Exception as e:
                            # If XML parsing fails, check for other image indicators
                            if 'drawing' in run.element.xml.lower() or 'image' in run.element.xml.lower():
                                page.insert_text((margin, y_position), "[Image]", fontsize=10)
                                y_position += 20

                # Start new page if needed
                if y_position > page.rect.height - 100:
                    page = pdf_doc.new_page()
                    y_position = 50

            pdf_doc.save(output_path)

        elif output_format == 'xlsx':
            wb = Workbook()
            ws = wb.active
            for i, paragraph in enumerate(doc.paragraphs):
                ws.cell(row=i+1, column=1, value=paragraph.text)
            wb.save(output_path)

        elif output_format == 'pptx':
            prs = Presentation()
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
                    tf = txBox.text_frame
                    tf.text = paragraph.text
            prs.save(output_path)

    # Excel to other formats
    elif input_ext == '.xlsx':
        wb = load_workbook(input_path)
        ws = wb.active

        if output_format == 'pdf':
            pdf_doc = fitz.open()
            page = pdf_doc.new_page()
            text = "\n".join([f"{cell.value}" for row in ws.rows for cell in row if cell.value])
            page.insert_text((50, 50), text)
            pdf_doc.save(output_path)

        elif output_format == 'docx':
            doc = Document()
            for row in ws.rows:
                text = " ".join([str(cell.value) for cell in row if cell.value])
                if text.strip():
                    doc.add_paragraph(text)
            doc.save(output_path)

        elif output_format == 'pptx':
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
            tf = txBox.text_frame
            text = "\n".join([f"{cell.value}" for row in ws.rows for cell in row if cell.value])
            tf.text = text
            prs.save(output_path)

    # PowerPoint to other formats
    elif input_ext == '.pptx':
        prs = Presentation(input_path)

        if output_format == 'pdf':
            pdf_doc = fitz.open()

            for slide_num, slide in enumerate(prs.slides):
                page = pdf_doc.new_page()
                page_width = page.rect.width
                page_height = page.rect.height

                # Get slide dimensions for scaling
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # Calculate scaling factors
                scale_x = page_width / slide_width
                scale_y = page_height / slide_height
                scale = min(scale_x, scale_y)  # Maintain aspect ratio

                for shape in slide.shapes:
                    # Get shape position and size
                    shape_left = shape.left * scale
                    shape_top = shape.top * scale
                    shape_width = shape.width * scale
                    shape_height = shape.height * scale

                    # Handle text shapes with positioning
                    if hasattr(shape, "text") and shape.text.strip():
                        # Calculate font size based on shape size
                        font_size = min(24, max(8, shape_height / 10))

                        # Position text according to original shape position
                        text_rect = fitz.Rect(shape_left, shape_top,
                                            shape_left + shape_width,
                                            shape_top + shape_height)

                        # Insert text with proper positioning
                        page.insert_text(text_rect.tl, shape.text, fontsize=font_size)

                    # Handle image shapes with proper positioning and sizing
                    elif hasattr(shape, 'image'):
                        try:
                            # Extract image from PowerPoint shape
                            image = shape.image
                            img_bytes = image.blob

                            # Position image according to original shape position and size
                            img_rect = fitz.Rect(shape_left, shape_top,
                                               shape_left + shape_width,
                                               shape_top + shape_height)

                            # Insert image into PDF with proper positioning
                            page.insert_image(img_rect, stream=img_bytes)

                        except Exception as e:
                            # If image extraction fails, add placeholder at correct position
                            page.insert_text((shape_left, shape_top), "[Image]", fontsize=10)
                            print(f"Warning: Could not extract image from slide {slide_num}: {e}")

                    # Handle other shape types that might contain images
                    elif shape.shape_type == 13:  # Picture shape type
                        page.insert_text((shape_left, shape_top), "[Image]", fontsize=10)

            pdf_doc.save(output_path)

        elif output_format == 'docx':
            doc = Document()

            for slide_num, slide in enumerate(prs.slides):
                # Add slide separator
                if slide_num > 0:
                    doc.add_page_break()

                doc.add_heading(f'Slide {slide_num + 1}', level=2)

                # Sort shapes by their vertical position to maintain layout order
                sorted_shapes = sorted(slide.shapes, key=lambda s: s.top)

                for shape in sorted_shapes:
                    # Handle text shapes with alignment detection
                    if hasattr(shape, "text") and shape.text.strip():
                        paragraph = doc.add_paragraph(shape.text)

                        # Determine alignment based on shape position
                        slide_width = prs.slide_width
                        shape_center = shape.left + (shape.width / 2)

                        from docx.enum.text import WD_ALIGN_PARAGRAPH
                        if shape_center < slide_width * 0.3:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        elif shape_center > slide_width * 0.7:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        elif shape_center > slide_width * 0.3 and shape_center < slide_width * 0.7:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

                    # Handle image shapes with proper sizing
                    elif hasattr(shape, 'image'):
                        try:
                            # Extract image from PowerPoint shape
                            image = shape.image
                            img_bytes = image.blob

                            # Create temporary image file
                            temp_img_path = f"temp_ppt_img_{slide_num}_{shape.shape_id}.png"
                            with open(temp_img_path, "wb") as img_file:
                                img_file.write(img_bytes)

                            # Calculate appropriate size based on original shape dimensions
                            from docx.shared import Inches

                            # Convert PowerPoint units to inches (PowerPoint uses EMUs)
                            original_width_inches = shape.width / 914400  # EMU to inches
                            original_height_inches = shape.height / 914400

                            # Scale to fit page (max 6.5 inches width for standard document)
                            max_width = 6.5
                            if original_width_inches > max_width:
                                scale_factor = max_width / original_width_inches
                                img_width = Inches(max_width)
                                img_height = Inches(original_height_inches * scale_factor)
                            else:
                                img_width = Inches(original_width_inches)
                                img_height = Inches(original_height_inches)

                            # Add image to Word document with proper sizing
                            paragraph = doc.add_paragraph()
                            run = paragraph.add_run()
                            run.add_picture(temp_img_path, width=img_width, height=img_height)

                            # Apply alignment based on image position in slide
                            shape_center = shape.left + (shape.width / 2)
                            from docx.enum.text import WD_ALIGN_PARAGRAPH
                            if shape_center < slide_width * 0.3:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                            elif shape_center > slide_width * 0.7:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                            else:
                                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

                            # Clean up temporary file
                            os.remove(temp_img_path)

                        except Exception as e:
                            # If image extraction fails, add placeholder
                            doc.add_paragraph("[Image could not be extracted]")
                            print(f"Warning: Could not extract image from slide {slide_num}: {e}")

                    # Handle other shape types that might contain images
                    elif shape.shape_type == 13:  # Picture shape type
                        doc.add_paragraph("[Image]")

            doc.save(output_path)

        elif output_format == 'xlsx':
            wb = Workbook()
            ws = wb.active
            row_num = 1
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        ws.cell(row=row_num, column=1, value=shape.text)
                        row_num += 1
            wb.save(output_path)

    # Image to PDF
    elif input_ext.lower() in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp']:
        if output_format == 'pdf':
            image = Image.open(input_path)
            pdf_doc = fitz.open()
            img_bytes = io.BytesIO()
            image.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            page = pdf_doc.new_page()
            page.insert_image(page.rect, stream=img_bytes.getvalue())
            pdf_doc.save(output_path)

    return True


def trim_media(file_path, start_time, end_time):
    supported_formats = {
        'audio': {'mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a'},
        'video': {'mp4', 'mkv', 'avi', 'mov', 'webm', 'flv'}
    }

    ext = os.path.splitext(file_path)[1][1:].lower()
    media_type = None
    for category, exts in supported_formats.items():
        if ext in exts:
            media_type = category
            break

    if media_type not in ['audio', 'video']:
        return False

    try:
        start = parse_time(start_time)
        end = parse_time(end_time)
    except ValueError:
        return False

    base, ext = os.path.splitext(file_path)
    output_path = f"{base}_trimmed{ext}"

    cmd = [
        ffmpeg_path, '-y',
        '-i', file_path,
        '-ss', str(start),
        '-to', str(end),
    ]

    if media_type == 'video':
        cmd += [
            '-c:v', 'libx264', '-preset', 'fast', '-crf', '23',
            '-c:a', 'aac', '-b:a', '192k'
        ]
    elif media_type == 'audio':
        codec_map = {
            'mp3': 'libmp3lame',
            'wav': 'pcm_s16le',
            'aac': 'aac',
            'flac': 'flac',
            'ogg': 'libvorbis',
            'm4a': 'aac'
        }
        cmd += ['-c:a', codec_map.get(ext[1:], 'copy')]

    cmd.append(output_path)

    try:
        subprocess.run(cmd, check=True, capture_output=True)
        return True
    except subprocess.CalledProcessError:
        return False
def get_available_formats(url):
    with YoutubeDL({'quiet': True}) as ydl:
        try:
            info = ydl.extract_info(url, download=False)
            formats = []
            for fmt in info.get('formats', []):
                if fmt.get('vcodec') != 'none':
                    res = fmt.get('resolution', 'unknown')
                    fps = fmt.get('fps', '?')
                    ext = fmt.get('ext', '?')
                    format_id = fmt.get('format_id', '')
                    size = fmt.get('filesize') or fmt.get('filesize_approx')
                    size_mb = f"{size/(1024*1024):.1f}MB" if size else "unknown size"
                    formats.append({
                        'format_id': format_id,
                        'resolution': res,
                        'fps': fps,
                        'ext': ext,
                        'size': size_mb,
                        'display': f"{res} {fps}fps | {ext.upper()} | {size_mb}"
                    })
            return formats
        except Exception as e:
            print(f"Error getting formats: {str(e)}")
            return []

class MediaUtilityGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("Media Utility")
        self.root.geometry("800x600")
        self.video_quality = None
        self.current_thread = None
        self.cancel_requested = False
        import queue
        self._result_queue = queue.Queue()

        # Initialize notebook
        self.notebook = ttk.Notebook(root)
        self.notebook.pack(fill='both', expand=True, padx=10, pady=5)

        # Create all tabs
        self.download_tab = ttk.Frame(self.notebook)
        self.convert_tab = ttk.Frame(self.notebook)
        self.batch_convert_tab = ttk.Frame(self.notebook)
        self.trim_tab = ttk.Frame(self.notebook)
        self.document_tab = ttk.Frame(self.notebook)  # Add document tab

        # Add all tabs to notebook
        self.notebook.add(self.download_tab, text='Download Media')
        self.notebook.add(self.convert_tab, text='Convert Media')
        self.notebook.add(self.batch_convert_tab, text='Batch Convert')
        self.notebook.add(self.trim_tab, text='Trim Media')
        self.notebook.add(self.document_tab, text='Document Convert')

        # Setup all tabs
        self.setup_download_tab()
        self.setup_convert_tab()
        self.setup_batch_convert_tab()
        self.setup_trim_tab()
        self.setup_document_tab()  # Setup document tab

        # Setup status frame
        self.status_frame = ttk.Frame(root)
        self.status_frame.pack(fill='x', padx=10, pady=5)
        self.progress_frame = ttk.Frame(self.status_frame)
        self.progress_frame.pack(fill='x', pady=5)
        self.progress = ttk.Progressbar(self.progress_frame, mode='indeterminate')
        self.progress.pack(side='left', fill='x', expand=True, padx=(0, 5))
        self.cancel_button = ttk.Button(self.progress_frame, text="Cancel",
                                      command=self.cancel_operation, state='disabled')
        self.cancel_button.pack(side='right')
        self.status_label = ttk.Label(self.status_frame, text="Ready")
        self.status_label.pack(pady=5)

    def setup_download_tab(self):
        url_frame = ttk.Frame(self.download_tab)
        url_frame.pack(fill='x', padx=10, pady=5)

        ttk.Label(url_frame, text="Enter URL:").pack(side='left', padx=5)
        self.url_entry = ttk.Entry(url_frame, width=60)
        self.url_entry.pack(side='left', padx=5, expand=True, fill='x')
        ttk.Button(url_frame, text="Check Available Formats",
                  command=self.check_formats).pack(side='left', padx=5)
        quality_frame = ttk.LabelFrame(self.download_tab, text="Video Quality")
        quality_frame.pack(pady=10, padx=10, fill='both', expand=True)

        self.quality_container = ttk.Frame(quality_frame)
        self.quality_container.pack(fill='both', expand=True)

        scrollbar = ttk.Scrollbar(self.quality_container)
        scrollbar.pack(side='right', fill='y')

        self.quality_listbox = tk.Listbox(self.quality_container,
                                        yscrollcommand=scrollbar.set,
                                        height=6)
        self.quality_listbox.pack(side='left', fill='both', expand=True)
        scrollbar.config(command=self.quality_listbox.yview)

        media_frame = ttk.LabelFrame(self.download_tab, text="Media Type")
        media_frame.pack(pady=10, padx=10, fill='x')

        self.media_type = tk.StringVar(value="video")
        ttk.Radiobutton(media_frame, text="Video", variable=self.media_type,
                       value="video", command=self.update_format_visibility).pack(side='left', padx=20)
        ttk.Radiobutton(media_frame, text="Audio", variable=self.media_type,
                       value="audio", command=self.update_format_visibility).pack(side='left', padx=20)

        self.audio_frame = ttk.LabelFrame(self.download_tab, text="Audio Format")
        self.audio_frame.pack(pady=10, padx=10, fill='x')
        self.download_audio_format = tk.StringVar(value="mp3")
        formats = ['mp3', 'aac', 'flac', 'wav', 'opus', 'm4a']
        for fmt in formats:
            ttk.Radiobutton(self.audio_frame, text=fmt.upper(), variable=self.download_audio_format, value=fmt).pack(side='left', padx=10)

        time_frame = ttk.LabelFrame(self.download_tab, text="Time Range (Optional)")
        time_frame.pack(pady=10, padx=10, fill='x')

        ttk.Label(time_frame, text="Start Time:").pack(side='left', padx=5)
        self.start_time = ttk.Entry(time_frame, width=10)
        self.start_time.pack(side='left', padx=5)

        ttk.Label(time_frame, text="End Time:").pack(side='left', padx=5)
        self.end_time = ttk.Entry(time_frame, width=10)
        self.end_time.pack(side='left', padx=5)
        # Add download location frame
        location_frame = ttk.LabelFrame(self.download_tab, text="Download Location (Optional)")
        location_frame.pack(pady=10, padx=10, fill='x')

        self.download_location = ttk.Entry(location_frame, width=50)
        self.download_location.pack(side='left', padx=5, fill='x', expand=True)

        ttk.Button(location_frame, text="Browse",
                  command=self.browse_download_location).pack(side='left', padx=5)

        ttk.Button(self.download_tab, text="Download",
                  command=self.start_download).pack(pady=20)
    def cancel_operation(self):
        if self.current_thread and self.current_thread.is_alive():
            self.cancel_requested = True
            self.update_status("Cancelling operation...")
            self.cancel_button.config(state='disabled')

    def start_progress(self):
        self.progress.start(10)
        self.cancel_button.config(state='normal')
        self.cancel_requested = False

    def stop_progress(self):
        self.progress.stop()
        self.cancel_button.config(state='disabled')
        self.current_thread = None

    def update_status(self, message, is_error=False):
        self.status_label.config(text=message,
                               foreground='red' if is_error else 'black')

    def _poll_result(self, on_success_msg="Done.", on_failure_msg="Operation failed!"):
        """Poll the result queue from the main thread via root.after()."""
        import queue
        try:
            status, payload = self._result_queue.get_nowait()
        except queue.Empty:
            self.root.after(100, self._poll_result, on_success_msg, on_failure_msg)
            return
        self.stop_progress()
        if status == "ok":
            self.update_status(payload if payload else on_success_msg)
        elif status == "cancelled":
            self.update_status("Operation cancelled.")
        else:  # "error"
            self.update_status(payload if payload else on_failure_msg, is_error=True)

    def check_formats(self):
        url = self.url_entry.get().strip()
        if not url:
            messagebox.showerror("Error", "Please enter a URL")
            return

        self.quality_listbox.delete(0, tk.END)
        self.quality_listbox.insert(tk.END, "Checking available formats...")
        self.root.update()

        def check_formats_thread():
            try:
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return
                formats = get_available_formats(url)
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                else:
                    self.root.after(0, self.update_quality_list, formats)
                    self._result_queue.put(("ok", "Formats loaded."))
            except Exception as e:
                self._result_queue.put(("error", f"Error checking formats: {str(e)}"))

        self.current_thread = threading.Thread(target=check_formats_thread, daemon=True)
        self.start_progress()
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Formats loaded.")


    def update_quality_list(self, formats):
        self.quality_listbox.delete(0, tk.END)
        self.quality_listbox.insert(tk.END, "Best Quality (Automatic)")
        self.formats = formats
        for fmt in formats:
            self.quality_listbox.insert(tk.END, fmt['display'])
    def update_format_visibility(self):
        if self.media_type.get() == 'audio':
            self.quality_container.pack_forget()
            self.audio_frame.pack(pady=10, padx=10, fill='x')
        else:
            self.quality_container.pack(fill='both', expand=True)
            self.audio_frame.pack(pady=10, padx=10, fill='x')

    def setup_convert_tab(self):
        # File Selection
        file_frame = ttk.Frame(self.convert_tab)
        file_frame.pack(pady=10, fill='x')

        self.convert_path = ttk.Entry(file_frame, width=60)
        self.convert_path.pack(side='left', padx=5)

        ttk.Button(file_frame, text="Browse",
                command=lambda: self.browse_file(self.convert_path)).pack(side='left')

        # Media Type Display
        self.media_type_label = ttk.Label(self.convert_tab, text="Media Type: None")
        self.media_type_label.pack(pady=5)

        # Format Selection Frame
        self.format_notebook = ttk.Notebook(self.convert_tab)
        self.format_notebook.pack(pady=10, padx=10, fill='both', expand=True)

        # Create tabs for each media type
        self.audio_frame = ttk.Frame(self.format_notebook)
        self.video_frame = ttk.Frame(self.format_notebook)
        self.image_frame = ttk.Frame(self.format_notebook)

        self.format_notebook.add(self.audio_frame, text='Audio Formats')
        self.format_notebook.add(self.video_frame, text='Video Formats')
        self.format_notebook.add(self.image_frame, text='Image Formats')

        # Audio Formats
        self.audio_format = tk.StringVar()
        audio_formats = ['mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a']
        self.setup_format_grid(self.audio_frame, audio_formats, self.audio_format)

        # Video Formats
        self.video_format = tk.StringVar()
        video_formats = ['mp4', 'mkv', 'avi', 'mov', 'webm', 'flv']
        self.setup_format_grid(self.video_frame, video_formats, self.video_format)

        # Image Formats
        self.image_format = tk.StringVar()
        image_formats = ['jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif']
        self.setup_format_grid(self.image_frame, image_formats, self.image_format)

        # Convert Button
        self.convert_button = ttk.Button(self.convert_tab, text="Convert",
                                    command=self.start_conversion)
        self.convert_button.pack(pady=20)

        # Bind file selection to update media type
        self.convert_path.bind('<KeyRelease>', self.update_media_type)
    def setup_format_grid(self, parent, formats, var):
        row = 0
        col = 0
        for fmt in formats:
            ttk.Radiobutton(parent, text=fmt.upper(),
                        variable=var, value=fmt).grid(
                            row=row, column=col, padx=10, pady=5)
            col += 1
            if col > 3:
                col = 0
                row += 1
    def update_media_type(self, event=None):
        file_path = self.convert_path.get()
        if not file_path or not os.path.exists(file_path):
            self.media_type_label.config(text="Media Type: None")
            return

        ext = os.path.splitext(file_path)[1][1:].lower()
        supported_formats = {
            'audio': {'mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a'},
            'video': {'mp4', 'mkv', 'avi', 'mov', 'webm', 'flv'},
            'image': {'jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif'}
        }

        media_type = None
        for category, exts in supported_formats.items():
            if ext in exts:
                media_type = category
                break

        if media_type:
            self.media_type_label.config(text=f"Media Type: {media_type.capitalize()}")
            tab_index = {'audio': 0, 'video': 1, 'image': 2}.get(media_type, 0)
            self.format_notebook.select(tab_index)
        else:
            self.media_type_label.config(text="Media Type: Unsupported format")
    def get_current_format_var(self):
        current_tab = self.format_notebook.select()
        tab_index = self.format_notebook.index(current_tab)
        return {
            0: self.audio_format,
            1: self.video_format,
            2: self.image_format
        }.get(tab_index)

    def setup_batch_convert_tab(self):
        file_frame = ttk.Frame(self.batch_convert_tab)
        file_frame.pack(pady=10, fill='x')

        self.batch_files = ttk.Entry(file_frame, width=60)
        self.batch_files.pack(side='left', padx=5)

        ttk.Button(file_frame, text="Browse",
                command=self.browse_multiple_files).pack(side='left')

        self.files_frame = ttk.LabelFrame(self.batch_convert_tab, text="Selected Files")
        self.files_frame.pack(pady=10, padx=10, fill='both', expand=True)

        self.files_text = tk.Text(self.files_frame, height=5, width=50)
        scrollbar = ttk.Scrollbar(self.files_frame, command=self.files_text.yview)
        self.files_text.configure(yscrollcommand=scrollbar.set)
        self.files_text.pack(side='left', fill='both', expand=True)
        scrollbar.pack(side='right', fill='y')
        self.files_text.config(state='disabled')

        self.batch_type_label = ttk.Label(self.batch_convert_tab, text="Media Type: None")
        self.batch_type_label.pack(pady=5)

        self.batch_format_notebook = ttk.Notebook(self.batch_convert_tab)
        self.batch_format_notebook.pack(pady=10, padx=10, fill='both', expand=True)

        self.batch_audio_frame = ttk.Frame(self.batch_format_notebook)
        self.batch_video_frame = ttk.Frame(self.batch_format_notebook)
        self.batch_image_frame = ttk.Frame(self.batch_format_notebook)

        self.batch_format_notebook.add(self.batch_audio_frame, text='Audio Formats')
        self.batch_format_notebook.add(self.batch_video_frame, text='Video Formats')
        self.batch_format_notebook.add(self.batch_image_frame, text='Image Formats')

        # Audio Formats
        self.batch_audio_format = tk.StringVar()
        audio_formats = ['mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a']
        self.setup_format_grid(self.batch_audio_frame, audio_formats, self.batch_audio_format)

        # Video Formats
        self.batch_video_format = tk.StringVar()
        video_formats = ['mp4', 'mkv', 'avi', 'mov', 'webm', 'flv']
        self.setup_format_grid(self.batch_video_frame, video_formats, self.batch_video_format)

        # Image Formats
        self.batch_image_format = tk.StringVar()
        image_formats = ['jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif']
        self.setup_format_grid(self.batch_image_frame, image_formats, self.batch_image_format)

        self.batch_convert_button = ttk.Button(self.batch_convert_tab, text="Convert All",
                                            command=self.start_batch_conversion)
        self.batch_convert_button.pack(pady=20)
    def browse_multiple_files(self):
        filenames = filedialog.askopenfilenames()
        if filenames:
            self.batch_files.delete(0, tk.END)
            self.batch_files.insert(0, ';'.join(filenames))

            self.files_text.config(state='normal')
            self.files_text.delete(1.0, tk.END)
            for file in filenames:
                self.files_text.insert(tk.END, f"{os.path.basename(file)}\n")
            self.files_text.config(state='disabled')

            self.update_batch_media_type(filenames)
    def update_batch_media_type(self, filenames):
        if not filenames:
            self.batch_type_label.config(text="Media Type: None")
            return

        supported_formats = {
            'audio': {'mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a'},
            'video': {'mp4', 'mkv', 'avi', 'mov', 'webm', 'flv'},
            'image': {'jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif'}
        }

        extensions = {os.path.splitext(f)[1][1:].lower() for f in filenames}

        media_types = set()
        for ext in extensions:
            for media_type, formats in supported_formats.items():
                if ext in formats:
                    media_types.add(media_type)
                    break

        if len(media_types) > 1:
            self.batch_type_label.config(text="Media Type: Mixed (not supported)")
            messagebox.showwarning("Warning",
                                "Mixed media types detected. Please select files of the same type.")
        elif len(media_types) == 1:
            media_type = media_types.pop()
            self.batch_type_label.config(text=f"Media Type: {media_type.capitalize()}")
            tab_index = {'audio': 0, 'video': 1, 'image': 2}.get(media_type, 0)
            self.batch_format_notebook.select(tab_index)
        else:
            self.batch_type_label.config(text="Media Type: Unsupported format")

    def get_current_batch_format_var(self):
        current_tab = self.batch_format_notebook.select()
        tab_index = self.batch_format_notebook.index(current_tab)
        return {
            0: self.convert_audio_format,
            1: self.video_format,
            2: self.image_format
        }.get(tab_index)
    def setup_trim_tab(self):
        file_frame = ttk.Frame(self.trim_tab)
        file_frame.pack(pady=10, fill='x')

        self.trim_path = ttk.Entry(file_frame, width=60)
        self.trim_path.pack(side='left', padx=5)

        ttk.Button(file_frame, text="Browse",
                  command=lambda: self.browse_file(self.trim_path)).pack(side='left')

        time_frame = ttk.LabelFrame(self.trim_tab, text="Time Range")
        time_frame.pack(pady=10, padx=10, fill='x')

        ttk.Label(time_frame, text="Start Time:").pack(side='left', padx=5)
        self.trim_start = ttk.Entry(time_frame, width=10)
        self.trim_start.pack(side='left', padx=5)

        ttk.Label(time_frame, text="End Time:").pack(side='left', padx=5)
        self.trim_end = ttk.Entry(time_frame, width=10)
        self.trim_end.pack(side='left', padx=5)

        ttk.Button(self.trim_tab, text="Trim Media",
                  command=self.start_trim).pack(pady=20)

    def setup_document_tab(self):
        # File Selection
        file_frame = ttk.Frame(self.document_tab)
        file_frame.pack(pady=10, fill='x')

        self.doc_path = ttk.Entry(file_frame, width=60)
        self.doc_path.pack(side='left', padx=5)

        ttk.Button(file_frame, text="Browse",
                command=lambda: self.browse_file(self.doc_path)).pack(side='left')

        # Format Selection
        format_frame = ttk.LabelFrame(self.document_tab, text="Target Format")
        format_frame.pack(pady=10, padx=10, fill='x')

        self.doc_format = tk.StringVar()
        formats = ['pdf', 'docx', 'xlsx', 'pptx']
        for fmt in formats:
            ttk.Radiobutton(format_frame, text=fmt.upper(),
                        variable=self.doc_format, value=fmt).pack(side='left', padx=10)

        # Conversion quality disclaimer
        warning_label = ttk.Label(
            self.document_tab,
            text="Note: Complex layouts may not convert perfectly. Best results with text-heavy documents.",
            foreground="gray",
            wraplength=500
        )
        warning_label.pack(pady=(0, 5), padx=10)

        # Convert Button
        ttk.Button(self.document_tab, text="Convert",
                command=self.start_doc_conversion).pack(pady=20)

    def start_doc_conversion(self):
        file_path = self.doc_path.get()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("Error", "Please select a valid file")
            return

        if not self.doc_format.get():
            messagebox.showerror("Error", "Please select a target format")
            return

        def doc_convert_thread():
            try:
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return
                success = convert_document(file_path, self.doc_format.get())
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                elif success:
                    self._result_queue.put(("ok", "Document conversion completed!"))
                else:
                    self._result_queue.put(("error", "Document conversion failed!"))
            except Exception as e:
                self._result_queue.put(("error", f"Error: {str(e)}"))

        self.current_thread = threading.Thread(target=doc_convert_thread, daemon=True)
        self.start_progress()
        self.update_status("Converting document...")
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Document conversion complete!")

    def browse_file(self, entry_widget):
        filename = filedialog.askopenfilename()
        if filename:
            entry_widget.delete(0, tk.END)
            entry_widget.insert(0, filename)

    def browse_download_location(self):
        directory = filedialog.askdirectory()
        if directory:
            self.download_location.delete(0, tk.END)
            self.download_location.insert(0, directory)

    def browse_multiple_files(self):
        filenames = filedialog.askopenfilenames()
        if filenames:
            self.batch_files.delete(0, tk.END)
            self.batch_files.insert(0, ';'.join(filenames))

    def start_download(self):
        url = self.url_entry.get().strip()
        if not url:
            messagebox.showerror("Error", "Please enter a URL")
            return

        quality = None
        if self.media_type.get() == 'video':
            selected_idx = self.quality_listbox.curselection()
            if selected_idx:
                if selected_idx[0] == 0:  # Best Quality
                    quality = 'bestvideo+bestaudio/best'
                elif hasattr(self, 'formats') and self.formats:
                    fmt = self.formats[selected_idx[0] - 1]
                    quality = f"{fmt['format_id']}+bestaudio/best"

        def download_thread():
            try:
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return
                success = download_media(
                    url=url,
                    platform=get_platform(url),
                    media_type=self.media_type.get(),
                    quality=quality,
                    start_time=self.start_time.get() if self.start_time.get() else None,
                    end_time=self.end_time.get() if self.end_time.get() else None,
                    audio_format=self.download_audio_format.get(),
                    output_dir=self.download_location.get() if self.download_location.get() else None,
                    video_codec="libx264",  # Target codec when conversion is needed
                    force_codec=False  # Only re-encode when codec is not already h264
                )
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                elif success:
                    out = self.download_location.get() if self.download_location.get() else os.getcwd()
                    self._result_queue.put(("ok", f"Download completed to {out}!"))
                else:
                    self._result_queue.put(("error", "Download failed!"))
            except Exception as e:
                self._result_queue.put(("error", f"Error: {str(e)}"))

        self.current_thread = threading.Thread(target=download_thread, daemon=True)
        self.start_progress()
        self.update_status("Downloading...")
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Download completed!")

    def start_batch_conversion(self):
        files = self.batch_files.get().split(';')
        if not files or not files[0]:
            messagebox.showerror("Error", "Please select files to convert")
            return

        format_var = self.get_current_batch_format_var()
        if not format_var or not format_var.get():
            messagebox.showerror("Error", "Please select a target format")
            return

        target_format = format_var.get()

        def batch_convert_thread():
            try:
                for i, file_path in enumerate(files):
                    if self.cancel_requested:
                        self._result_queue.put(("cancelled", None))
                        return
                    # Intermediate status updates are safe via root.after from the worker
                    self.root.after(0, self.update_status, f"Converting file {i+1} of {len(files)}...")
                    success = convert_images([file_path], target_format)
                    if not success and not self.cancel_requested:
                        self._result_queue.put(("error", f"Failed to convert {os.path.basename(file_path)}"))
                        return
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                else:
                    self._result_queue.put(("ok", f"Successfully converted {len(files)} files!"))
            except Exception as e:
                self._result_queue.put(("error", f"Error: {str(e)}"))

        self.current_thread = threading.Thread(target=batch_convert_thread, daemon=True)
        self.start_progress()
        self.update_status("Converting files...")
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Batch conversion complete!")



    def start_conversion(self):
        file_path = self.convert_path.get()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("Error", "Please select a valid file")
            return

        format_var = self.get_current_format_var()
        if not format_var or not format_var.get():
            messagebox.showerror("Error", "Please select a target format")
            return

        target_format = format_var.get()

        def convert_thread():
            try:
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return

                ext = os.path.splitext(file_path)[1][1:].lower()
                if ext == target_format:
                    self._result_queue.put(("error", "Source and target formats are the same!"))
                    return

                supported_formats = {
                    'audio': {'mp3', 'wav', 'aac', 'flac', 'ogg', 'm4a'},
                    'video': {'mp4', 'mkv', 'avi', 'mov', 'webm', 'flv'},
                    'image': {'jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp', 'heic', 'heif'}
                }

                media_type = None
                for category, formats in supported_formats.items():
                    if ext in formats:
                        media_type = category
                        break

                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return

                if media_type == "image":
                    success = convert_images([file_path], target_format)
                else:
                    base = os.path.splitext(file_path)[0]
                    output_path = f"{base}_converted.{target_format}"
                    if media_type == 'video' and target_format in supported_formats['audio']:
                        cmd = [
                            ffmpeg_path, '-y',
                            '-i', file_path,
                            '-vn',
                            '-acodec', 'libmp3lame' if target_format == 'mp3' else target_format,
                            output_path
                        ]
                    else:
                        cmd = [ffmpeg_path, '-y', '-i', file_path, output_path]
                    if self.cancel_requested:
                        self._result_queue.put(("cancelled", None))
                        return
                    result = subprocess.run(cmd, capture_output=True)
                    success = result.returncode == 0

                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                elif success:
                    self._result_queue.put(("ok", "Conversion completed successfully!"))
                else:
                    self._result_queue.put(("error", "Conversion failed!"))

            except Exception as e:
                self._result_queue.put(("error", f"Error: {str(e)}"))

        self.current_thread = threading.Thread(target=convert_thread, daemon=True)
        self.start_progress()
        self.update_status("Converting...")
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Conversion complete!")


    def start_trim(self):
        file_path = self.trim_path.get()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("Error", "Please select a valid file")
            return

        def trim_thread():
            try:
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                    return
                success = trim_media(
                    file_path,
                    self.trim_start.get(),
                    self.trim_end.get()
                )
                if self.cancel_requested:
                    self._result_queue.put(("cancelled", None))
                elif success:
                    self._result_queue.put(("ok", "Trimming completed successfully!"))
                else:
                    self._result_queue.put(("error", "Trimming failed!"))
            except Exception as e:
                self._result_queue.put(("error", f"Error: {str(e)}"))

        self.current_thread = threading.Thread(target=trim_thread, daemon=True)
        self.start_progress()
        self.update_status("Trimming media...")
        self.current_thread.start()
        self.root.after(100, self._poll_result, "Trimming complete!")

def main():
    import queue as _queue

    root = tk.Tk()
    root.withdraw()  # Hide main window until dependency check completes

    # Splash window shown while checking/installing dependencies
    splash = tk.Toplevel(root)
    splash.title("Starting Media Utility")
    splash.geometry("380x110")
    splash.resizable(False, False)
    ttk.Label(splash, text="Checking dependencies, please wait...", padding=(20, 15)).pack()
    _bar = ttk.Progressbar(splash, mode='indeterminate')
    _bar.pack(fill='x', padx=20, pady=(0, 15))
    _bar.start(10)

    _result_q = _queue.Queue()
    threading.Thread(target=lambda: _result_q.put(check_dependencies()), daemon=True).start()

    def _poll_deps():
        try:
            dep_error = _result_q.get_nowait()
        except _queue.Empty:
            root.after(100, _poll_deps)
            return
        _bar.stop()
        splash.destroy()
        root.deiconify()
        root._gui = MediaUtilityGUI(root)  # attach to root so it stays alive for root.mainloop()
        if dep_error == "ffmpeg_missing":
            messagebox.showerror(
                "Missing Dependency",
                "FFmpeg was not found on this system.\n\n"
                "Media conversion, trimming, and download features will not work.\n\n"
                "Install FFmpeg and add it to your system PATH, or place ffmpeg.exe "
                "in the same directory as this application.\n\n"
                "Download from: https://ffmpeg.org/download.html"
            )

    root.after(100, _poll_deps)
    root.mainloop()

if __name__ == "__main__":
    main()
